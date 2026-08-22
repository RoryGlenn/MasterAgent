"""Deterministic local draft and preview generators for Phase 3."""

from __future__ import annotations

import difflib
import fnmatch
import hashlib
import json
import os
import re
import stat
from collections.abc import Mapping, Sequence
from dataclasses import dataclass
from email.message import EmailMessage
from io import BytesIO
from pathlib import Path
from threading import Lock
from typing import Any

from master_agent.directory_safety import PinnedDirectory, pin_directory
from master_agent.errors import ConfigurationError, ConnectorError
from master_agent.models import (
    ActionState,
    AgentAction,
    ExecutionResult,
    ResourceRef,
    VerificationResult,
)
from master_agent.platform_runtime import require_persistent_state_platform
from master_agent.resource_limits import (
    MAX_LOCAL_ARTIFACT_BYTES,
    MAX_RUN_ARTIFACT_BYTES,
)

_SAFE_NAME = re.compile(r"[^A-Za-z0-9._-]+")


@dataclass(frozen=True, slots=True)
class GeneratedArtifact:
    """Metadata for a locally generated draft."""

    path: Path
    sha256: str
    media_type: str
    size: int

    def to_dict(self) -> dict[str, Any]:
        """Return JSON-compatible artifact metadata."""

        return {
            "path": str(self.path),
            "sha256": self.sha256,
            "media_type": self.media_type,
            "size": self.size,
        }


class ArtifactBudget:
    """One thread-safe aggregate byte budget shared by a complete local run."""

    def __init__(self, max_bytes: int = MAX_RUN_ARTIFACT_BYTES) -> None:
        if (
            isinstance(max_bytes, bool)
            or not isinstance(max_bytes, int)
            or max_bytes <= 0
            or max_bytes > MAX_RUN_ARTIFACT_BYTES
        ):
            raise ValueError(
                f"artifact budget must be between 1 and {MAX_RUN_ARTIFACT_BYTES} bytes"
            )
        self._max_bytes = max_bytes
        self._used_bytes = 0
        self._lock = Lock()

    @property
    def max_bytes(self) -> int:
        """Return the immutable aggregate ceiling."""

        return self._max_bytes

    @property
    def used_bytes(self) -> int:
        """Return the bytes retained by successful artifact publications."""

        with self._lock:
            return self._used_bytes

    def reserve(self, amount: int) -> None:
        """Reserve bytes before any final artifact name is created."""

        if isinstance(amount, bool) or not isinstance(amount, int) or amount < 0:
            raise ValueError("artifact reservation must be a non-negative integer")
        with self._lock:
            if self._used_bytes + amount > self._max_bytes:
                raise ConnectorError(
                    "aggregate generated artifact budget would be exceeded"
                )
            self._used_bytes += amount

    def release(self, amount: int) -> None:
        """Release a failed transaction's exact prior reservation."""

        with self._lock:
            if (
                isinstance(amount, bool)
                or not isinstance(amount, int)
                or amount < 0
                or amount > self._used_bytes
            ):
                raise RuntimeError("artifact budget release is inconsistent")
            self._used_bytes -= amount


class _LocalDraftConnector:
    """Base connector for generation under one controlled output root."""

    def __init__(
        self,
        *,
        system: str,
        capabilities: frozenset[str],
        output_root: Path | PinnedDirectory,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        require_persistent_state_platform()
        self._system = system
        self._capabilities = capabilities
        self._output_directory = pin_directory(output_root)
        self._output_root = self._output_directory.path
        self._artifact_budget = artifact_budget or ArtifactBudget()
        supplied_limits = dict(output_limits or {})
        self._output_limits: dict[str, int] = {}
        for capability in capabilities:
            limit = supplied_limits.get(capability, MAX_LOCAL_ARTIFACT_BYTES)
            if (
                isinstance(limit, bool)
                or not isinstance(limit, int)
                or limit <= 0
                or limit > MAX_LOCAL_ARTIFACT_BYTES
            ):
                raise ConfigurationError(
                    f"local artifact output quota is invalid for {capability}"
                )
            self._output_limits[capability] = limit

    @property
    def system(self) -> str:
        """Return connector system."""

        return self._system

    @property
    def capabilities(self) -> frozenset[str]:
        """Return supported local-generation capabilities."""

        return self._capabilities

    def close(self) -> None:
        """Release the connector-owned output-directory pin."""

        self._output_directory.close()

    def read(self, resource: ResourceRef) -> dict[str, object] | None:
        """Read generated artifact metadata by resource identifier."""

        self._output_directory.validate()
        matches = sorted(
            name
            for name in os.listdir(self._output_directory.fileno())
            if fnmatch.fnmatchcase(
                name,
                f"{_safe_name(resource.resource_id)}.*",
            )
        )
        if not matches:
            return None
        path = self._output_root / matches[0]
        digest, size = _inspect_artifact(
            self._output_directory,
            path,
            max_bytes=MAX_LOCAL_ARTIFACT_BYTES,
        )
        return {
            "path": str(path),
            "sha256": digest,
            "size": size,
        }

    def verify(
        self,
        action: AgentAction,
        result: ExecutionResult,
    ) -> VerificationResult:
        """Verify the artifact exists and matches its recorded digest."""

        after = result.after or {}
        path_value = after.get("path")
        digest = after.get("sha256")
        if not isinstance(path_value, str) or not isinstance(digest, str):
            return VerificationResult(
                action_id=action.action_id,
                verified=False,
                observed=None,
                message="generated artifact metadata was incomplete",
            )
        path = Path(path_value)
        try:
            observed_digest, observed_size = _inspect_artifact(
                self._output_directory,
                path,
                max_bytes=self._output_limits[action.capability],
            )
        except ConnectorError:
            observed_digest = None
            observed_size = None
        verified = observed_digest is not None and observed_digest == digest
        observed = {
            "path": str(path),
            "sha256": observed_digest,
            "size": observed_size,
        }
        return VerificationResult(
            action_id=action.action_id,
            verified=verified,
            observed=observed,
            message=(
                "verified local artifact digest"
                if verified
                else "local artifact digest mismatch"
            ),
        )

    def _artifact_path(
        self,
        action: AgentAction,
        *,
        default_suffix: str,
    ) -> Path:
        raw = str(action.parameters.get("output_name", "")).strip()
        if raw:
            name = Path(raw).name
            if name != raw or raw in {".", ".."}:
                raise ConnectorError("output_name must be a single filename")
            name = _safe_name(name)
        else:
            name = f"{_safe_name(action.target.resource_id)}{default_suffix}"
        if not name:
            raise ConnectorError("generated artifact filename is empty")
        try:
            self._output_directory.validate()
        except ConfigurationError as error:
            raise ConnectorError("generated artifact destination changed") from error
        return self._output_root / name

    def _write_json(
        self,
        action: AgentAction,
        path: Path,
        payload: Mapping[str, Any],
    ) -> GeneratedArtifact:
        return _write_json(
            self._output_directory,
            path,
            payload,
            artifact_budget=self._artifact_budget,
            max_output_bytes=self._output_limits[action.capability],
        )

    def _write_text(
        self,
        action: AgentAction,
        path: Path,
        text: str,
        media_type: str,
    ) -> GeneratedArtifact:
        return _write_text(
            self._output_directory,
            path,
            text,
            media_type,
            artifact_budget=self._artifact_budget,
            max_output_bytes=self._output_limits[action.capability],
        )

    def _write_bytes(
        self,
        action: AgentAction,
        path: Path,
        payload: bytes | memoryview,
        media_type: str,
    ) -> GeneratedArtifact:
        return _write_bytes(
            self._output_directory,
            path,
            payload,
            media_type,
            artifact_budget=self._artifact_budget,
            max_output_bytes=self._output_limits[action.capability],
        )

    def _write_bundle(
        self,
        action: AgentAction,
        files: Sequence[tuple[Path, bytes | memoryview, str]],
    ) -> tuple[GeneratedArtifact, ...]:
        """Create a companion-file bundle with transaction-owned rollback."""

        return _write_bundle(
            self._output_directory,
            files,
            artifact_budget=self._artifact_budget,
            max_output_bytes=self._output_limits[action.capability],
        )

    def _result(
        self,
        action: AgentAction,
        artifact: GeneratedArtifact,
        *,
        message: str,
        extra: Mapping[str, Any] | None = None,
    ) -> ExecutionResult:
        after = artifact.to_dict()
        if extra:
            after.update(dict(extra))
        return ExecutionResult(
            action_id=action.action_id,
            state=ActionState.SUCCEEDED,
            before=None,
            after=after,
            connector_reference=str(artifact.path),
            message=message,
        )


class JiraDraftConnector(_LocalDraftConnector):
    """Generate Jira update/comment/transition proposals without publishing."""

    _CAPABILITIES = frozenset(
        {
            "jira.issue.update.draft",
            "jira.issue.comment.draft",
            "jira.issue.transition.draft",
        }
    )

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="jira",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Write a JSON proposal with a Markdown preview."""

        _require_capability(action, self)
        before = _mapping(action.parameters.get("before"), "before", default={})
        if action.capability == "jira.issue.update.draft":
            change = _mapping(action.parameters.get("fields"), "fields")
            after = {**before, **change}
            operation = "update"
        elif action.capability == "jira.issue.comment.draft":
            body = _required_text(action.parameters, "body")
            after = {"comment": body}
            operation = "comment"
        else:
            transition_id = _required_text(action.parameters, "transition_id")
            after = {
                **before,
                "transition_id": transition_id,
                "target_status": str(action.parameters.get("target_status", "")).strip()
                or None,
            }
            operation = "transition"

        payload = {
            "schema": "master-agent/jira-draft@1",
            "system": "jira",
            "issue": action.target.resource_id,
            "operation": operation,
            "expected_version": action.target.expected_version,
            "before": before,
            "after": after,
            "publish": False,
        }
        path = self._artifact_path(action, default_suffix=".jira-draft.json")
        preview_path = path.with_suffix(".md")
        preview = _render_preview(
            title=f"Jira {operation} draft — {action.target.resource_id}",
            before=before,
            after=after,
        )
        artifact, preview_artifact = self._write_bundle(
            action,
            (
                (path, _json_bytes(payload), "application/json"),
                (preview_path, preview.encode("utf-8"), "text/markdown"),
            ),
        )
        return self._result(
            action,
            artifact,
            message="generated Jira draft without publishing",
            extra={"preview": preview_artifact.to_dict(), "operation": operation},
        )


class ConfluenceDraftConnector(_LocalDraftConnector):
    """Generate Confluence create/update proposals without publishing."""

    _CAPABILITIES = frozenset(
        {"confluence.page.create.draft", "confluence.page.update.draft"}
    )

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="confluence",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Write Confluence page proposal and before/after preview."""

        _require_capability(action, self)
        before = _mapping(action.parameters.get("before"), "before", default={})
        title = _required_text(action.parameters, "title")
        body = _required_text(action.parameters, "body")
        representation = str(action.parameters.get("representation", "storage")).strip()
        if representation not in {"storage", "atlas_doc_format"}:
            raise ConnectorError(
                "Confluence representation must be storage or atlas_doc_format"
            )
        after = {
            "title": title,
            "body": body,
            "representation": representation,
            "space_id": action.parameters.get("space_id"),
            "parent_id": action.parameters.get("parent_id"),
        }
        operation = "create" if action.capability.endswith("create.draft") else "update"
        payload = {
            "schema": "master-agent/confluence-draft@1",
            "system": "confluence",
            "page_id": action.target.resource_id,
            "operation": operation,
            "expected_version": action.target.expected_version,
            "before": before,
            "after": after,
            "publish": False,
        }
        path = self._artifact_path(action, default_suffix=".confluence-draft.json")
        preview = _render_preview(
            title=f"Confluence {operation} draft — {title}",
            before=before,
            after=after,
        )
        artifact, preview_artifact = self._write_bundle(
            action,
            (
                (path, _json_bytes(payload), "application/json"),
                (path.with_suffix(".md"), preview.encode("utf-8"), "text/markdown"),
            ),
        )
        return self._result(
            action,
            artifact,
            message="generated Confluence draft without publishing",
            extra={"preview": preview_artifact.to_dict(), "operation": operation},
        )


class OutlookDraftConnector(_LocalDraftConnector):
    """Generate RFC 5322 email drafts without sending."""

    _CAPABILITIES = frozenset({"outlook.email.draft"})

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="outlook",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Create a local ``.eml`` file and metadata manifest."""

        _require_capability(action, self)
        to = _string_list(action.parameters.get("to"), "to", required=True)
        cc = _string_list(action.parameters.get("cc"), "cc")
        bcc = _string_list(action.parameters.get("bcc"), "bcc")
        subject = _required_text(action.parameters, "subject")
        body = _required_text(action.parameters, "body")
        message = EmailMessage()
        message["To"] = ", ".join(to)
        if cc:
            message["Cc"] = ", ".join(cc)
        if bcc:
            message["Bcc"] = ", ".join(bcc)
        message["Subject"] = subject
        message.set_content(body)
        path = self._artifact_path(action, default_suffix=".eml")
        metadata = {
            "schema": "master-agent/outlook-draft@1",
            "to": list(to),
            "cc": list(cc),
            "bcc": list(bcc),
            "subject": subject,
            "body_sha256": hashlib.sha256(body.encode("utf-8")).hexdigest(),
            "send": False,
        }
        artifact, manifest = self._write_bundle(
            action,
            (
                (path, message.as_bytes(), "message/rfc822"),
                (path.with_suffix(".json"), _json_bytes(metadata), "application/json"),
            ),
        )
        return self._result(
            action,
            artifact,
            message="generated local Outlook email draft without sending",
            extra={
                "manifest": manifest.to_dict(),
                "recipients": len(to) + len(cc) + len(bcc),
            },
        )


class TeamsDraftConnector(_LocalDraftConnector):
    """Generate Teams message drafts without posting."""

    _CAPABILITIES = frozenset({"teams.message.draft"})

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="teams",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Write a Teams draft as Markdown plus structured metadata."""

        _require_capability(action, self)
        body = _required_text(action.parameters, "body")
        recipient_type = str(action.parameters.get("recipient_type", "chat")).strip()
        if recipient_type not in {"chat", "channel", "user", "team"}:
            raise ConnectorError("recipient_type must be chat, channel, user, or team")
        recipient_id = _required_text(action.parameters, "recipient_id")
        path = self._artifact_path(action, default_suffix=".teams-draft.md")
        header = (
            f"# Teams message draft\n\n"
            f"- Recipient type: `{recipient_type}`\n"
            f"- Recipient ID: `{recipient_id}`\n"
            f"- Posted: **no**\n\n"
            f"## Message\n\n{body}\n"
        )
        metadata = {
            "schema": "master-agent/teams-draft@1",
            "recipient_type": recipient_type,
            "recipient_id": recipient_id,
            "body_sha256": hashlib.sha256(body.encode("utf-8")).hexdigest(),
            "post": False,
        }
        artifact, manifest = self._write_bundle(
            action,
            (
                (path, header.encode("utf-8"), "text/markdown"),
                (path.with_suffix(".json"), _json_bytes(metadata), "application/json"),
            ),
        )
        return self._result(
            action,
            artifact,
            message="generated local Teams message draft without posting",
            extra={"manifest": manifest.to_dict()},
        )


class PowerPointDraftConnector(_LocalDraftConnector):
    """Generate local PowerPoint presentations from typed slide data."""

    _CAPABILITIES = frozenset({"powerpoint.presentation.generate"})

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="powerpoint",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Render a bounded presentation with no external publishing."""

        _require_capability(action, self)
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as error:  # pragma: no cover - declared dependency.
            raise ConnectorError(
                "python-pptx is required for PowerPoint generation"
            ) from error

        title = _required_text(action.parameters, "title")
        slides_value = action.parameters.get("slides")
        if slides_value is None:
            sections = action.parameters.get("sections", [])
            if not isinstance(sections, Sequence) or isinstance(sections, (str, bytes)):
                raise ConnectorError("slides or sections must be a list")
            slides_value = [
                {"title": str(section), "bullets": []} for section in sections
            ]
        if not isinstance(slides_value, Sequence) or isinstance(
            slides_value, (str, bytes)
        ):
            raise ConnectorError("slides must be a list")
        if len(slides_value) > 40:
            raise ConnectorError("PowerPoint generation is limited to 40 slides")

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = str(
                action.parameters.get("subtitle", "Generated by Master Agent")
            )

        for raw_slide in slides_value:
            if not isinstance(raw_slide, Mapping):
                raise ConnectorError("each slide must be an object")
            slide_title = str(raw_slide.get("title", "Untitled")).strip() or "Untitled"
            bullets = raw_slide.get("bullets", [])
            if not isinstance(bullets, Sequence) or isinstance(bullets, (str, bytes)):
                raise ConnectorError("slide bullets must be a list")
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = slide_title[:160]
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for index, bullet in enumerate(bullets[:12]):
                paragraph = (
                    text_frame.paragraphs[0]
                    if index == 0
                    else text_frame.add_paragraph()
                )
                paragraph.text = str(bullet)[:800]
                paragraph.level = 0
                paragraph.font.size = Pt(20)

        path = self._artifact_path(action, default_suffix=".pptx")
        if path.suffix.lower() != ".pptx":
            path = path.with_suffix(".pptx")
        output = BytesIO()
        presentation.save(output)
        payload = output.getbuffer()
        try:
            artifact = self._write_bytes(
                action,
                path,
                payload,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        finally:
            payload.release()
        return self._result(
            action,
            artifact,
            message="generated local PowerPoint presentation",
            extra={"slide_count": len(presentation.slides)},
        )


class RepositoryDraftConnector(_LocalDraftConnector):
    """Generate source-code patches and branch plans without touching a repo."""

    _CAPABILITIES = frozenset({"repository.patch.generate", "repository.branch.plan"})

    def __init__(
        self,
        output_root: Path | PinnedDirectory,
        *,
        artifact_budget: ArtifactBudget | None = None,
        output_limits: Mapping[str, int] | None = None,
    ) -> None:
        super().__init__(
            system="repository",
            capabilities=self._CAPABILITIES,
            output_root=output_root,
            artifact_budget=artifact_budget,
            output_limits=output_limits,
        )

    def execute(self, action: AgentAction) -> ExecutionResult:
        """Generate a unified diff or branch-operation plan."""

        _require_capability(action, self)
        if action.capability == "repository.patch.generate":
            relative_path = _safe_relative_path(
                _required_text(action.parameters, "relative_path")
            )
            before = str(action.parameters.get("before_text", ""))
            after = str(action.parameters.get("after_text", ""))
            diff = "".join(
                difflib.unified_diff(
                    before.splitlines(keepends=True),
                    after.splitlines(keepends=True),
                    fromfile=f"a/{relative_path}",
                    tofile=f"b/{relative_path}",
                )
            )
            if not diff:
                raise ConnectorError("patch generation produced no changes")
            path = self._artifact_path(action, default_suffix=".patch")
            artifact = self._write_text(action, path, diff, "text/x-diff")
            return self._result(
                action,
                artifact,
                message="generated local source-code patch",
                extra={"relative_path": relative_path},
            )

        branch = _required_text(action.parameters, "branch")
        base = _required_text(action.parameters, "base")
        _validate_branch(branch)
        _validate_branch(base)
        payload = {
            "schema": "master-agent/repository-branch-plan@1",
            "branch": branch,
            "base": base,
            "remote": str(action.parameters.get("remote", "origin")).strip()
            or "origin",
            "push": False,
        }
        path = self._artifact_path(action, default_suffix=".branch-plan.json")
        artifact = self._write_json(action, path, payload)
        return self._result(
            action,
            artifact,
            message="generated local repository branch plan",
        )


def _require_capability(action: AgentAction, connector: _LocalDraftConnector) -> None:
    if action.target.system != connector.system:
        raise ConnectorError(
            f"connector {connector.system} cannot execute target {action.target.system}"
        )
    if action.capability not in connector.capabilities:
        raise ConnectorError(f"unsupported draft capability: {action.capability}")


def _render_preview(
    *,
    title: str,
    before: Mapping[str, Any],
    after: Mapping[str, Any],
) -> str:
    before_text = json.dumps(before, indent=2, sort_keys=True, ensure_ascii=False)
    after_text = json.dumps(after, indent=2, sort_keys=True, ensure_ascii=False)
    diff = "".join(
        difflib.unified_diff(
            before_text.splitlines(keepends=True),
            after_text.splitlines(keepends=True),
            fromfile="before.json",
            tofile="after.json",
        )
    )
    return (
        f"# {title}\n\n"
        "This is a local proposal. It has not been published.\n\n"
        "## Diff\n\n```diff\n"
        f"{diff}```\n"
    )


def _write_json(
    directory: PinnedDirectory,
    path: Path,
    payload: Mapping[str, Any],
    *,
    artifact_budget: ArtifactBudget,
    max_output_bytes: int,
) -> GeneratedArtifact:
    return _write_bytes(
        directory,
        path,
        _json_bytes(payload),
        "application/json",
        artifact_budget=artifact_budget,
        max_output_bytes=max_output_bytes,
    )


def _write_text(
    directory: PinnedDirectory,
    path: Path,
    text: str,
    media_type: str,
    *,
    artifact_budget: ArtifactBudget,
    max_output_bytes: int,
) -> GeneratedArtifact:
    return _write_bytes(
        directory,
        path,
        text.encode("utf-8"),
        media_type,
        artifact_budget=artifact_budget,
        max_output_bytes=max_output_bytes,
    )


def _write_bytes(
    directory: PinnedDirectory,
    path: Path,
    payload: bytes | memoryview,
    media_type: str,
    *,
    artifact_budget: ArtifactBudget,
    max_output_bytes: int,
) -> GeneratedArtifact:
    return _write_bundle(
        directory,
        ((path, payload, media_type),),
        artifact_budget=artifact_budget,
        max_output_bytes=max_output_bytes,
    )[0]


def _write_bundle(
    directory: PinnedDirectory,
    files: Sequence[tuple[Path, bytes | memoryview, str]],
    *,
    artifact_budget: ArtifactBudget,
    max_output_bytes: int,
) -> tuple[GeneratedArtifact, ...]:
    """Create all files or remove only earlier transaction-owned files."""

    if not files:
        raise ConnectorError("generated artifact bundle must not be empty")
    if (
        isinstance(max_output_bytes, bool)
        or not isinstance(max_output_bytes, int)
        or max_output_bytes <= 0
        or max_output_bytes > MAX_LOCAL_ARTIFACT_BYTES
    ):
        raise ConnectorError("generated artifact output quota is invalid")
    total_bytes = sum(len(payload) for _path, payload, _media_type in files)
    if total_bytes > max_output_bytes:
        raise ConnectorError(
            "generated artifact bundle exceeds the capability output quota"
        )
    descriptor = directory.fileno()
    artifact_budget.reserve(total_bytes)
    owned: list[tuple[str, tuple[int, int, int, int, int], GeneratedArtifact]] = []
    try:
        for path, payload, media_type in files:
            name = _pinned_name(directory, path)
            artifact, identity = _create_bytes(
                directory,
                name,
                path,
                payload,
                media_type,
            )
            owned.append((name, identity, artifact))
        return tuple(item[2] for item in owned)
    except BaseException:
        for name, identity, _ in reversed(owned):
            _unlink_if_identity(descriptor, name, identity)
        try:
            os.fsync(descriptor)
        except OSError:
            pass
        artifact_budget.release(total_bytes)
        raise


def write_artifact_bundle(
    output_root: Path | PinnedDirectory,
    files: Sequence[tuple[Path, bytes | memoryview, str]],
    *,
    artifact_budget: ArtifactBudget | None = None,
    max_output_bytes: int = MAX_LOCAL_ARTIFACT_BYTES,
) -> tuple[GeneratedArtifact, ...]:
    """Create a local artifact bundle through one pinned output directory.

    Every final name is created with ``O_EXCL``. If a later companion cannot
    be created, only earlier files whose exact identities belong to this
    transaction are removed.
    """

    require_persistent_state_platform()
    with pin_directory(output_root) as directory:
        return _write_bundle(
            directory,
            files,
            artifact_budget=artifact_budget or ArtifactBudget(),
            max_output_bytes=max_output_bytes,
        )


def _create_bytes(
    directory: PinnedDirectory,
    name: str,
    path: Path,
    payload: bytes | memoryview,
    media_type: str,
) -> tuple[GeneratedArtifact, tuple[int, int, int, int, int]]:
    """Create and verify one final-name file without overwrite or rename."""

    descriptor = directory.fileno()
    file_descriptor = -1
    created_identity: tuple[int, int, int, int, int] | None = None
    completed = False
    expected_size = len(payload)
    expected_digest = hashlib.sha256(payload).hexdigest()
    try:
        file_descriptor = os.open(
            name,
            os.O_RDWR | os.O_CREAT | os.O_EXCL | getattr(os, "O_NOFOLLOW", 0),
            0o600,
            dir_fd=descriptor,
        )
        os.fchmod(file_descriptor, 0o600)
        created_identity = _restricted_identity(os.fstat(file_descriptor))
        remaining = memoryview(payload)
        while remaining:
            written = os.write(file_descriptor, remaining)
            if written <= 0:
                raise OSError("short generated artifact write")
            remaining = remaining[written:]
        os.fsync(file_descriptor)
        directory.validate()
        published = os.stat(name, dir_fd=descriptor, follow_symlinks=False)
        if _restricted_identity(published) != created_identity:
            raise ConfigurationError("generated artifact publication was replaced")
        os.lseek(file_descriptor, 0, os.SEEK_SET)
        observed_digest = hashlib.sha256()
        observed_size = 0
        remaining_bytes = expected_size
        while remaining_bytes:
            chunk = os.read(file_descriptor, min(1024 * 1024, remaining_bytes))
            if not chunk:
                raise ConfigurationError(
                    "generated artifact bytes changed during write"
                )
            observed_size += len(chunk)
            remaining_bytes -= len(chunk)
            observed_digest.update(chunk)
        if os.read(file_descriptor, 1):
            raise ConfigurationError("generated artifact bytes changed during write")
        if (
            observed_size != expected_size
            or observed_digest.hexdigest() != expected_digest
        ):
            raise ConfigurationError("generated artifact bytes changed during write")
        os.fsync(descriptor)
        directory.validate()
        completed = True
    except FileExistsError as error:
        raise ConnectorError(
            "generated artifact already exists; use a fresh output name or directory"
        ) from error
    except (OSError, ConfigurationError) as error:
        raise ConnectorError("generated artifact destination changed") from error
    finally:
        if file_descriptor >= 0:
            os.close(file_descriptor)
        if not completed and created_identity is not None:
            _unlink_if_identity(descriptor, name, created_identity)
    if created_identity is None:  # pragma: no cover - successful open invariant.
        raise ConnectorError("generated artifact identity is missing")
    return (
        GeneratedArtifact(
            path=path,
            sha256=expected_digest,
            media_type=media_type,
            size=expected_size,
        ),
        created_identity,
    )


def _json_bytes(payload: Mapping[str, Any]) -> bytes:
    """Render deterministic pretty JSON bytes for draft bundles."""

    return (
        json.dumps(payload, indent=2, ensure_ascii=False, sort_keys=True) + "\n"
    ).encode("utf-8")


def _inspect_artifact(
    directory: PinnedDirectory,
    path: Path,
    *,
    max_bytes: int,
) -> tuple[str, int]:
    """Stream a bounded artifact digest without retaining a second full copy."""

    name = _pinned_name(directory, path)
    descriptor = directory.fileno()
    try:
        file_descriptor = os.open(
            name,
            os.O_RDONLY | getattr(os, "O_NOFOLLOW", 0),
            dir_fd=descriptor,
        )
    except OSError as error:
        raise ConnectorError("generated artifact destination changed") from error
    try:
        metadata = os.fstat(file_descriptor)
        try:
            identity = _restricted_identity(metadata)
        except ConfigurationError as error:
            raise ConnectorError("generated artifact file is unsafe") from error
        if metadata.st_size > max_bytes:
            raise ConnectorError("generated artifact exceeds its output quota")
        digest = hashlib.sha256()
        total = 0
        remaining = metadata.st_size
        while remaining:
            chunk = os.read(file_descriptor, min(1024 * 1024, remaining))
            if not chunk:
                raise ConnectorError("generated artifact changed during verification")
            total += len(chunk)
            remaining -= len(chunk)
            digest.update(chunk)
        if os.read(file_descriptor, 1):
            raise ConnectorError("generated artifact changed during verification")
        if _restricted_identity(os.fstat(file_descriptor)) != identity:
            raise ConnectorError("generated artifact changed during verification")
        published = os.stat(name, dir_fd=descriptor, follow_symlinks=False)
        if _restricted_identity(published) != identity:
            raise ConnectorError("generated artifact changed during verification")
        directory.validate()
        return digest.hexdigest(), total
    finally:
        os.close(file_descriptor)


def _pinned_name(directory: PinnedDirectory, path: Path) -> str:
    if path.parent != directory.path or path.name in {"", ".", ".."}:
        raise ConnectorError("generated artifact escaped the output root")
    try:
        directory.validate()
    except ConfigurationError as error:
        raise ConnectorError("generated artifact destination changed") from error
    return path.name


def _restricted_identity(
    metadata: os.stat_result,
) -> tuple[int, int, int, int, int]:
    identity = (
        metadata.st_dev,
        metadata.st_ino,
        metadata.st_uid,
        stat.S_IMODE(metadata.st_mode),
        metadata.st_nlink,
    )
    if (
        not stat.S_ISREG(metadata.st_mode)
        or metadata.st_uid != os.getuid()
        or stat.S_IMODE(metadata.st_mode) != 0o600
        or metadata.st_nlink != 1
    ):
        raise ConfigurationError("generated artifact file is unsafe")
    return identity


def _unlink_if_identity(
    parent_descriptor: int,
    name: str,
    expected: tuple[int, int, int, int, int],
) -> None:
    try:
        current = os.stat(name, dir_fd=parent_descriptor, follow_symlinks=False)
    except FileNotFoundError:
        return
    try:
        if _restricted_identity(current) != expected:
            return
        os.unlink(name, dir_fd=parent_descriptor)
    except (ConfigurationError, FileNotFoundError):
        return


def _safe_name(value: str) -> str:
    rendered = _SAFE_NAME.sub("-", value.strip()).strip(".-")
    return rendered[:180] or "artifact"


def _mapping(
    value: Any, name: str, *, default: Mapping[str, Any] | None = None
) -> dict[str, Any]:
    if value is None and default is not None:
        return dict(default)
    if not isinstance(value, Mapping):
        raise ConnectorError(f"parameter must be an object: {name}")
    return dict(value)


def _required_text(parameters: Mapping[str, Any], key: str) -> str:
    value = str(parameters.get(key, "")).strip()
    if not value:
        raise ConnectorError(f"missing required parameter: {key}")
    return value


def _string_list(value: Any, name: str, *, required: bool = False) -> tuple[str, ...]:
    if value is None:
        result: tuple[str, ...] = ()
    elif isinstance(value, str):
        result = tuple(item.strip() for item in value.split(",") if item.strip())
    elif isinstance(value, Sequence) and not isinstance(value, (str, bytes)):
        result = tuple(str(item).strip() for item in value if str(item).strip())
    else:
        raise ConnectorError(f"parameter must be a string list: {name}")
    if required and not result:
        raise ConnectorError(f"parameter must not be empty: {name}")
    return result


def _safe_relative_path(value: str) -> str:
    path = Path(value)
    if path.is_absolute() or ".." in path.parts or not path.parts:
        raise ConnectorError("relative_path must remain inside the repository")
    return path.as_posix()


def _validate_branch(value: str) -> None:
    if (
        not value
        or value.startswith("-")
        or ".." in value
        or value.endswith("/")
        or any(character.isspace() for character in value)
        or any(token in value for token in ("~", "^", ":", "?", "*", "[", "\\"))
    ):
        raise ConnectorError("unsafe Git branch name")
