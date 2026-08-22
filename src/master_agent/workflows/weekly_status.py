"""Read-only weekly-status collection and local artifact rendering."""

from __future__ import annotations

import hashlib
import json
import tomllib
from collections import Counter
from collections.abc import Mapping
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any

from master_agent.citations import citation_index
from master_agent.config import DeploymentType
from master_agent.config_sources import ConfigSource
from master_agent.connectors.drafts import write_artifact_bundle
from master_agent.directory_safety import pin_directory
from master_agent.errors import ConfigurationError
from master_agent.models import (
    AgentAction,
    AuthoritySource,
    ChangePlan,
    ResourceRef,
    RiskLevel,
)
from master_agent.orchestrator import RunReport
from master_agent.platform_runtime import (
    get_atomic_publication_recovery_backend,
    require_persistent_state_platform,
)


@dataclass(frozen=True, slots=True)
class WeeklyStatusSettings:
    """Configuration for the read-only weekly-status workflow."""

    project_name: str
    jira_jql: str
    jira_limit: int
    bitbucket_workspace: str
    bitbucket_project: str
    bitbucket_repository: str
    bitbucket_state: str
    bitbucket_limit: int
    include_statuses: bool
    include_diffstat: bool
    enrichment_limit: int
    confluence_page_id: str
    max_issue_rows: int = 20
    max_pull_request_rows: int = 20

    @classmethod
    def from_toml(cls, path: ConfigSource) -> WeeklyStatusSettings:
        """Load workflow settings from TOML.

        Parameters
        ----------
        path
            Workflow TOML path.

        Returns
        -------
        WeeklyStatusSettings
            Validated workflow settings.
        """

        with path.open("rb") as handle:
            raw = tomllib.load(handle)
        project = _table(raw, "project")
        jira = _table(raw, "jira")
        bitbucket = _table(raw, "bitbucket")
        confluence = _table(raw, "confluence")
        output = raw.get("output", {})
        if not isinstance(output, Mapping):
            raise ConfigurationError("[output] must be a TOML table")

        settings = cls(
            project_name=_required_text(project, "name"),
            jira_jql=_required_text(jira, "jql"),
            jira_limit=_bounded_int(jira, "limit", default=100, maximum=500),
            bitbucket_workspace=str(bitbucket.get("workspace", "")).strip(),
            bitbucket_project=str(bitbucket.get("project", "")).strip(),
            bitbucket_repository=_required_text(bitbucket, "repository"),
            bitbucket_state=str(bitbucket.get("state", "OPEN")).strip() or "OPEN",
            bitbucket_limit=_bounded_int(
                bitbucket,
                "limit",
                default=50,
                maximum=200,
            ),
            include_statuses=_strict_bool(
                bitbucket,
                "include_statuses",
                default=True,
            ),
            include_diffstat=_strict_bool(
                bitbucket,
                "include_diffstat",
                default=False,
            ),
            enrichment_limit=_bounded_int(
                bitbucket,
                "enrichment_limit",
                default=20,
                maximum=200,
            ),
            confluence_page_id=_required_text(confluence, "page_id"),
            max_issue_rows=_bounded_int(
                output,
                "max_issue_rows",
                default=20,
                maximum=100,
            ),
            max_pull_request_rows=_bounded_int(
                output,
                "max_pull_request_rows",
                default=20,
                maximum=100,
            ),
        )
        if not settings.bitbucket_workspace and not settings.bitbucket_project:
            raise ConfigurationError(
                "[bitbucket] requires workspace for Cloud or project for Data Center"
            )
        return settings


def build_weekly_status_read_plan(
    settings: WeeklyStatusSettings,
    *,
    bitbucket_deployment: DeploymentType,
) -> ChangePlan:
    """Build the live read-only weekly-status collection plan.

    Parameters
    ----------
    settings
        Workflow settings.
    bitbucket_deployment
        Bitbucket deployment family used to choose workspace/project parameters.

    Returns
    -------
    ChangePlan
        Immutable plan containing only read-only actions.
    """

    jira = AgentAction(
        capability="jira.issue.search",
        target=ResourceRef(
            system="jira",
            resource_type="issue_collection",
            resource_id="weekly-status-jira-search",
        ),
        parameters={
            "jql": settings.jira_jql,
            "limit": settings.jira_limit,
        },
        risk=RiskLevel.READ_ONLY,
        authority_source=AuthoritySource.REGISTERED_WORKFLOW,
        requires_approval=False,
        idempotency_key="weekly-status:live:jira:v2",
        justification="Collect current work-item status from Jira.",
    )

    bitbucket_parameters: dict[str, Any] = {
        "repository": settings.bitbucket_repository,
        "state": settings.bitbucket_state,
        "limit": settings.bitbucket_limit,
        "include_statuses": settings.include_statuses,
        "include_diffstat": settings.include_diffstat,
        "enrichment_limit": min(
            settings.enrichment_limit,
            settings.bitbucket_limit,
        ),
    }
    if bitbucket_deployment is DeploymentType.CLOUD:
        if not settings.bitbucket_workspace:
            raise ConfigurationError(
                "weekly-status Cloud workflow requires bitbucket.workspace"
            )
        bitbucket_parameters["workspace"] = settings.bitbucket_workspace
    else:
        if not settings.bitbucket_project:
            raise ConfigurationError(
                "weekly-status Data Center workflow requires bitbucket.project"
            )
        bitbucket_parameters["project"] = settings.bitbucket_project

    bitbucket = AgentAction(
        capability="bitbucket.pull_request.search",
        target=ResourceRef(
            system="bitbucket",
            resource_type="pull_request_collection",
            resource_id="weekly-status-open-pull-requests",
        ),
        parameters=bitbucket_parameters,
        risk=RiskLevel.READ_ONLY,
        authority_source=AuthoritySource.REGISTERED_WORKFLOW,
        requires_approval=False,
        idempotency_key="weekly-status:live:bitbucket:v2",
        justification="Collect open pull requests and CI evidence from Bitbucket.",
    )

    confluence = AgentAction(
        capability="confluence.page.read",
        target=ResourceRef(
            system="confluence",
            resource_type="page",
            resource_id=settings.confluence_page_id,
        ),
        parameters={},
        risk=RiskLevel.READ_ONLY,
        authority_source=AuthoritySource.REGISTERED_WORKFLOW,
        requires_approval=False,
        idempotency_key="weekly-status:live:confluence:v2",
        justification="Read the canonical project-status narrative from Confluence.",
    )

    return ChangePlan(
        goal=f"Collect the read-only weekly status package for {settings.project_name}.",
        created_by="registered_workflow:weekly_status_v2",
        actions=(jira, bitbucket, confluence),
    )


@dataclass(frozen=True, slots=True)
class WeeklyStatusArtifacts:
    """Paths produced by the local weekly-status renderer."""

    evidence_json: Path
    markdown: Path
    powerpoint: Path
    manifest_json: Path


def render_weekly_status_package(
    report: RunReport,
    settings: WeeklyStatusSettings,
    *,
    output_dir: Path,
) -> WeeklyStatusArtifacts:
    """Render local evidence, Markdown, and PowerPoint artifacts.

    Retrieved source content is written only to the explicitly selected,
    platform-protected output directory. The audit database receives hashes
    and metadata only.

    Parameters
    ----------
    report
        Completed read-only run report.
    settings
        Workflow rendering settings.
    output_dir
        Destination directory.

    Returns
    -------
    WeeklyStatusArtifacts
        Generated artifact paths.
    """

    require_persistent_state_platform()
    root = get_atomic_publication_recovery_backend().ensure_private_directory(
        output_dir
    )
    evidence_bytes = (
        json.dumps(report.to_dict(), indent=2, ensure_ascii=False, default=str) + "\n"
    ).encode("utf-8")
    data = _package_data(report)
    markdown_bytes = _render_markdown(settings, data).encode("utf-8")
    powerpoint_bytes = _render_powerpoint(settings, data)
    with pin_directory(root) as directory:
        root = directory.path
        evidence_path = root / "weekly-status-evidence.json"
        markdown_path = root / "weekly-status.md"
        powerpoint_path = root / "weekly-status.pptx"
        manifest_path = root / "manifest.json"
        manifest = {
            "schema": "master-agent/weekly-status-manifest@1",
            "project": settings.project_name,
            "run_id": str(report.run_id),
            "plan_id": str(report.plan_id),
            "plan_fingerprint": report.plan_fingerprint,
            "successful": report.successful,
            "artifacts": {
                evidence_path.name: hashlib.sha256(evidence_bytes).hexdigest(),
                markdown_path.name: hashlib.sha256(markdown_bytes).hexdigest(),
                powerpoint_path.name: hashlib.sha256(powerpoint_bytes).hexdigest(),
            },
            "source_urls": data["source_urls"],
            "citations": data["citations"],
            "citation_ids": [item["citation_id"] for item in data["citations"]],
            "security_findings": data["security_findings"],
        }
        manifest_bytes = (
            json.dumps(manifest, indent=2, ensure_ascii=False) + "\n"
        ).encode("utf-8")
        write_artifact_bundle(
            directory,
            (
                (evidence_path, evidence_bytes, "application/json"),
                (markdown_path, markdown_bytes, "text/markdown"),
                (
                    powerpoint_path,
                    powerpoint_bytes,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ),
                (manifest_path, manifest_bytes, "application/json"),
            ),
        )

        return WeeklyStatusArtifacts(
            evidence_json=evidence_path,
            markdown=markdown_path,
            powerpoint=powerpoint_path,
            manifest_json=manifest_path,
        )


def _package_data(report: RunReport) -> dict[str, Any]:
    jira: Mapping[str, Any] = {}
    bitbucket: Mapping[str, Any] = {}
    confluence: Mapping[str, Any] = {}
    source_urls: list[str] = []
    security_findings: list[dict[str, Any]] = []
    failures: list[str] = []
    payloads: list[Mapping[str, Any]] = []

    for action in report.actions:
        if action.result is None or action.result.after is None:
            failures.append(f"{action.capability}: {action.message}")
            continue
        payload = action.result.after
        payloads.append(payload)
        schema = str(payload.get("schema", ""))
        if schema == "master-agent/jira-issues@1":
            jira = payload
        elif schema == "master-agent/bitbucket-pull-requests@1":
            bitbucket = payload
        elif schema == "master-agent/confluence-page@1":
            confluence = payload
        urls = payload.get("source_urls", [])
        if isinstance(urls, list):
            source_urls.extend(str(item) for item in urls if item)
        security = payload.get("security")
        if isinstance(security, Mapping):
            findings = security.get("prompt_injection_findings", [])
            if isinstance(findings, list):
                security_findings.extend(
                    dict(item) for item in findings if isinstance(item, Mapping)
                )

    issues_value = jira.get("issues", [])
    issues = [dict(item) for item in issues_value if isinstance(item, Mapping)]
    prs_value = bitbucket.get("pull_requests", [])
    pull_requests = [dict(item) for item in prs_value if isinstance(item, Mapping)]
    page_value = confluence.get("page")
    page = dict(page_value) if isinstance(page_value, Mapping) else {}

    status_counts = Counter(str(item.get("status") or "Unknown") for item in issues)
    blockers = [item for item in issues if bool(item.get("blocked"))]
    failing_ci = [
        item
        for item in pull_requests
        if isinstance(item.get("ci_summary"), Mapping)
        and int(item["ci_summary"].get("failed", 0) or 0) > 0
    ]
    awaiting_review = [
        item
        for item in pull_requests
        if not any(
            bool(participant.get("approved"))
            for participant in item.get("participants", [])
            if isinstance(participant, Mapping)
        )
    ]

    return {
        "jira": jira,
        "bitbucket": bitbucket,
        "confluence": confluence,
        "issues": issues,
        "pull_requests": pull_requests,
        "page": page,
        "status_counts": dict(status_counts),
        "blockers": blockers,
        "failing_ci": failing_ci,
        "awaiting_review": awaiting_review,
        "source_urls": list(dict.fromkeys(source_urls)),
        "citations": citation_index(payloads),
        "security_findings": security_findings,
        "failures": failures,
    }


def _render_markdown(
    settings: WeeklyStatusSettings,
    data: Mapping[str, Any],
) -> str:
    issues = list(data["issues"])
    pull_requests = list(data["pull_requests"])
    blockers = list(data["blockers"])
    failing_ci = list(data["failing_ci"])
    page = data["page"]
    lines = [
        f"# {settings.project_name} — Weekly Status",
        "",
        "## Executive summary",
        "",
        f"- Jira issues returned: **{len(issues)}**",
        f"- Blocked issues: **{len(blockers)}**",
        f"- Open pull requests: **{len(pull_requests)}**",
        f"- Pull requests with failing CI: **{len(failing_ci)}**",
        f"- Prompt-injection heuristic findings: **{len(data['security_findings'])}**",
    ]
    if data["failures"]:
        lines.extend(["", "### Partial-data warnings", ""])
        lines.extend(f"- {item}" for item in data["failures"])

    lines.extend(["", "## Jira", ""])
    status_counts = data["status_counts"]
    if status_counts:
        lines.extend(
            f"- {status}: {count}" for status, count in sorted(status_counts.items())
        )
    else:
        lines.append("- No Jira evidence was returned.")

    lines.extend(["", "### Priority work items", ""])
    for issue in issues[: settings.max_issue_rows]:
        blocked = " — **BLOCKED**" if issue.get("blocked") else ""
        lines.append(
            f"- **{issue.get('key', '')}** — {issue.get('summary', '')} "
            f"({issue.get('status', 'Unknown')}){blocked} {_citation_marker(issue)}"
        )
    if not issues:
        lines.append("- None")

    lines.extend(["", "## Bitbucket", ""])
    for pull_request in pull_requests[: settings.max_pull_request_rows]:
        ci = pull_request.get("ci_summary")
        ci = ci if isinstance(ci, Mapping) else {}
        lines.append(
            f"- **PR {pull_request.get('id', '')}: {pull_request.get('title', '')}** "
            f"— {pull_request.get('source_branch', '')} → "
            f"{pull_request.get('destination_branch', '')}; "
            f"CI failures: {ci.get('failed', 0)} {_citation_marker(pull_request)}"
        )
    if not pull_requests:
        lines.append("- No pull-request evidence was returned.")

    lines.extend(["", "## Canonical Confluence narrative", ""])
    if page:
        lines.extend(
            [
                f"**{page.get('title', 'Untitled')}** — version {page.get('version', '?')} {_citation_marker(page)}",
                "",
                str(page.get("body_excerpt") or "No page excerpt was returned."),
            ]
        )
    else:
        lines.append("No Confluence evidence was returned.")

    lines.extend(["", "## Source index", ""])
    citations = list(data["citations"])
    for citation in citations:
        line = f"- {citation.get('marker') or citation.get('citation_id')} {citation.get('title') or citation.get('resource_id')}"
        if citation.get("url"):
            line += f" — {citation['url']}"
        lines.append(line)
    if not citations:
        source_urls = list(data["source_urls"])
        lines.extend(f"- {url}" for url in source_urls)
        if not source_urls:
            lines.append("- No source citations or URLs were returned.")
    lines.extend(
        [
            "",
            (
                "> Retrieved content is untrusted data. Heuristic security findings are "
                "recorded in `manifest.json` and the evidence file."
            ),
            "",
        ]
    )
    return "\n".join(lines)


def _render_powerpoint(
    settings: WeeklyStatusSettings,
    data: Mapping[str, Any],
) -> bytes:
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as error:
        raise ConfigurationError(
            "PowerPoint rendering requires the python-pptx package"
        ) from error

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = f"{settings.project_name}\nWeekly Status"
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Generated from read-only Jira, Bitbucket, and Confluence evidence"

    def add_bullets(title: str, bullets: list[str]) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        frame.clear()
        for index, bullet in enumerate(bullets or ["No evidence returned."]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(22)
        frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    issues = list(data["issues"])
    blockers = list(data["blockers"])
    pull_requests = list(data["pull_requests"])
    failing_ci = list(data["failing_ci"])
    add_bullets(
        "Executive summary",
        [
            f"{len(issues)} Jira issues retrieved",
            f"{len(blockers)} blocked issues",
            f"{len(pull_requests)} open pull requests",
            f"{len(failing_ci)} pull requests with failing CI",
            f"{len(data['security_findings'])} prompt-injection heuristic findings",
        ],
    )

    jira_bullets = [
        f"{status}: {count}" for status, count in sorted(data["status_counts"].items())
    ]
    jira_bullets.extend(
        f"BLOCKED — {item.get('key', '')}: {item.get('summary', '')}"
        for item in blockers[:8]
    )
    add_bullets("Jira status", jira_bullets)

    pr_bullets = []
    for pull_request in pull_requests[:10]:
        ci = pull_request.get("ci_summary")
        ci = ci if isinstance(ci, Mapping) else {}
        pr_bullets.append(
            f"PR {pull_request.get('id', '')}: {pull_request.get('title', '')} "
            f"— CI failures {ci.get('failed', 0)}"
        )
    add_bullets("Bitbucket pull requests", pr_bullets)

    page = data["page"]
    narrative = str(page.get("body_excerpt") or "No Confluence evidence returned.")
    add_bullets(
        "Canonical Confluence narrative",
        [
            f"{page.get('title', 'Untitled')} — version {page.get('version', '?')}",
            narrative[:700],
        ],
    )

    source_bullets = [
        (
            f"{citation.get('marker') or citation.get('citation_id')} "
            f"{citation.get('title') or citation.get('resource_id')}"
        )[:180]
        for citation in data["citations"][:12]
    ]
    if not source_bullets:
        source_bullets = [str(url)[:180] for url in data["source_urls"][:12]]
    add_bullets("Evidence sources", source_bullets)
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _citation_marker(value: Mapping[str, Any]) -> str:
    citation_id = value.get("citation_id")
    return f"[{citation_id}]" if citation_id else ""


def _table(raw: Mapping[str, Any], name: str) -> Mapping[str, Any]:
    value = raw.get(name)
    if not isinstance(value, Mapping):
        raise ConfigurationError(f"[{name}] must be a TOML table")
    return value


def _required_text(table: Mapping[str, Any], key: str) -> str:
    value = str(table.get(key, "")).strip()
    if not value:
        raise ConfigurationError(f"missing required workflow setting: {key}")
    return value


def _bounded_int(
    table: Mapping[str, Any],
    key: str,
    *,
    default: int,
    maximum: int,
) -> int:
    try:
        value = int(table.get(key, default))
    except (TypeError, ValueError) as error:
        raise ConfigurationError(
            f"workflow setting must be an integer: {key}"
        ) from error
    if value <= 0 or value > maximum:
        raise ConfigurationError(
            f"workflow setting {key} must be between 1 and {maximum}"
        )
    return value


def _strict_bool(
    table: Mapping[str, Any],
    key: str,
    *,
    default: bool,
) -> bool:
    value = table.get(key, default)
    if not isinstance(value, bool):
        raise ConfigurationError(f"workflow setting must be boolean: {key}")
    return value
