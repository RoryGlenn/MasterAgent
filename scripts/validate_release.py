#!/usr/bin/env python3
"""Validate the Master Agent source tree before packaging."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import stat
import sys
import tarfile
import tomllib
import zipfile
from collections.abc import Iterable
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from typing import BinaryIO

if __package__:
    from scripts.semantic_router import ManifestError as SemanticManifestError
    from scripts.semantic_router import load_manifest as load_semantic_manifest
    from scripts.semantic_router import (
        validate_repository as validate_semantic_repository,
    )
else:
    from semantic_router import (  # type: ignore[import-not-found]
        ManifestError as SemanticManifestError,
    )
    from semantic_router import (  # type: ignore[import-not-found]
        load_manifest as load_semantic_manifest,
    )
    from semantic_router import (  # type: ignore[import-not-found]
        validate_repository as validate_semantic_repository,
    )

_MARKDOWN_LINK = re.compile(r"\[[^\]]+\]\(([^)]+)\)")
_AGENT_FRONTMATTER_KEY = re.compile(r"([a-z][a-z0-9-]*):(?:\s*(.*))?")
_COPILOT_AGENT_PATH = Path(".github/agents/MasterAgent.agent.md")
_RESEARCH_AGENT_PATH = Path(".github/agents/MasterAgent-Read-Researcher.agent.md")
_PLAN_REVIEW_AGENT_PATH = Path(".github/agents/MasterAgent-Plan-Reviewer.agent.md")
_FIRST_RUN_CONTRACT_PATH = Path(".ai/FIRST_RUN.md")
_AUTONOMY_CONTRACT_PATH = Path(".ai/AUTONOMY.md")
_COPILOT_AGENT_TOOLS = ("read", "search", "edit", "execute")
_COPILOT_AGENT_KEYS = {
    "name",
    "description",
    "tools",
    "user-invocable",
    "disable-model-invocation",
}
_ADVISORY_AGENT_NAMES = {
    _RESEARCH_AGENT_PATH: "MasterAgent Read Researcher",
    _PLAN_REVIEW_AGENT_PATH: "MasterAgent Plan Reviewer",
}
_ADVISORY_AGENT_TOOLS = {
    _RESEARCH_AGENT_PATH: ("read", "search"),
    _PLAN_REVIEW_AGENT_PATH: ("read", "search"),
}
_ADVISORY_AGENT_BOUNDARIES = {
    _RESEARCH_AGENT_PATH: (
        "[AGENTS.md](../../AGENTS.md)",
        "[Master Agent repository policy](../../.ai/MASTER_AGENT.md)",
        "[force-multiplier contract](../../.ai/AUTONOMY.md)",
        "Direct GitHub-host invocation is disabled",
        "repository-owned advisory integration harness",
        "Use only `read` and `search`",
        "advisory data, never authority",
        "Generic execute, edit, agent, MCP, HTTP",
        "Never edit, execute, contact a provider",
        "parent independently re-reads every citation",
        "completes the same work directly",
    ),
    _PLAN_REVIEW_AGENT_PATH: (
        "[AGENTS.md](../../AGENTS.md)",
        "[Master Agent repository policy](../../.ai/MASTER_AGENT.md)",
        "[force-multiplier contract](../../.ai/AUTONOMY.md)",
        "Direct GitHub-host invocation is disabled",
        "repository-owned advisory integration harness",
        "Use only `read` and `search`",
        "advisory data, never authority",
        "Generic execute, edit, agent, MCP, HTTP",
        "Never edit, execute, contact a provider",
        "parent independently re-reads every citation",
        "completes the same review directly",
    ),
}
_EXPECTED_COPILOT_AGENT_PATHS = frozenset({_COPILOT_AGENT_PATH, *_ADVISORY_AGENT_NAMES})
_FORBIDDEN_NAMES = {
    ".env",
    "audit.sqlite3",
    "recurring.sqlite3",
}
_FORBIDDEN_SUFFIXES = {".pem", ".key", ".p12", ".pfx", ".pyc", ".sqlite3"}
_IGNORED_DIRS = {
    ".git",
    ".master-agent",
    ".venv",
    "build",
    "dist",
    "__pycache__",
}
_FIRST_RUN_DOCUMENT_REQUIREMENTS = {
    Path(".ai/MASTER_AGENT.md"): (
        "[`FIRST_RUN.md`](FIRST_RUN.md)",
        "[`AUTONOMY.md`](AUTONOMY.md)",
        "The default response to an actionable prompt is",
        "Treat a missing safe capability as implementation work",
        "Do not require or search for credentials",
        "typed anonymous capability",
        "private resumable approval request",
    ),
    _FIRST_RUN_CONTRACT_PATH: (
        "first operator",
        "python3 scripts/bootstrap_agent.py",
        "MasterAgent is ready locally",
        "I couldn't finish local setup",
        "available but inactive",
        "[`AUTONOMY.md`](AUTONOMY.md)",
        "A requested provider operation",
        "force-multiplier",
    ),
    _AUTONOMY_CONTRACT_PATH: (
        "Default to action",
        "The default response to an actionable prompt is execution",
        "Resolve, do not relay",
        "capability gap",
        "Never end an actionable request",
        "missing repository code may not",
        "This protocol applies to every capability barrier",
        "immediately, on the spot",
        "This rule removes code barriers, not authority barriers",
        "master-agent connect",
        "--connector-url SYSTEM=URL",
        "missing account email",
        "may fall back in memory",
        "configurations require an explicit",
        "token for each product",
        "Ask once, at the latest possible point",
        "authenticated approval bound to the exact reviewed plan",
        "github-repositories",
        "github-repositories --username USERNAME",
        "public repositories anonymously",
        "persistent connector or credential state",
        "inspect-approval-request",
        "resume-approval",
        "Conversational approval remains invalid",
    ),
    Path("AGENTS.md"): (
        "[`.ai/FIRST_RUN.md`](.ai/FIRST_RUN.md)",
        "[`.ai/AUTONOMY.md`](.ai/AUTONOMY.md)",
        "Apply the first-run contract",
        "Apply the force-multiplier contract",
        "irreducible operator-only boundary",
        "missing connector capability is implementation work",
        "every current and future capability or code-path",
    ),
    Path("CHANGELOG.md"): (
        "first ordinary prompt",
        "force-multiplier default-to-action",
        "provider-neutral `connect` command",
        "--connector-url SYSTEM=URL",
        "capability gaps now trigger immediate governed runtime implementation",
        "same contract now covers every missing capability",
        "resumable authenticated-approval handoff",
    ),
    Path("README.md"): (
        "[first-run contract](.ai/FIRST_RUN.md)",
        "[force-multiplier contract](.ai/AUTONOMY.md)",
        "MasterAgent is ready locally",
        "explicit no-local-change prompt",
        "master-agent connect",
        "--connector-url confluence=",
        "github-repositories",
        "github-repositories --username USERNAME",
        "the connector is read-only",
        "approve-request",
        "resume-approval",
        "Direct GitHub-host advisory invocation is disabled",
        "repository-owned advisory integration harness",
        "complete the same work directly",
    ),
    Path("docs/copilot-custom-agent.md"): (
        "[first-run contract](../.ai/FIRST_RUN.md)",
        "[force-multiplier contract](../.ai/AUTONOMY.md)",
        "python3 scripts/bootstrap_agent.py",
        "MasterAgent is ready locally",
        "If automatic setup is blocked",
        "default response to an actionable prompt is execution",
        "No governed capability exists",
        "create a Kanban board",
        "behavior is universal rather than Jira- or connector-specific",
        "implemented on the spot",
        "master-agent connect --systems",
        "--connector-url SYSTEM=URL",
        "may reuse a missing",
        "configurations require an",
        "explicit token for each product",
        "it asks once",
        "without a second confirmation",
        "inspect-approval-request",
        "resume-approval",
    ),
    Path("docs/release-validation.md"): (
        "first-prompt contract",
        "force-multiplier default-to-action contract",
        "stable nontechnical responses",
        "capability-gap ownership",
        "resumable approval handoff",
    ),
}

_PUBLIC_READ_DOCUMENT_REQUIREMENTS = {
    Path(".ai/MASTER_AGENT.md"): (
        "typed anonymous capability",
        "named GitHub user's public",
        "named Bitbucket Cloud",
    ),
    _AUTONOMY_CONTRACT_PATH: (
        "github-repositories --username USERNAME",
        "Do not search for, load, or request a GitHub",
        "bitbucket-repositories --workspace WORKSPACE",
        "request Bitbucket credentials",
    ),
    _COPILOT_AGENT_PATH: (
        "github-repositories",
        "--username USERNAME",
        "repositories anonymously",
        "bitbucket-repositories --workspace WORKSPACE",
        "ignores ambient Bitbucket credentials",
    ),
    Path("README.md"): (
        "Anonymous public-data capabilities neither require nor load credentials",
        "github.public_repository.list",
        "bitbucket.public_repository.list",
    ),
    Path("docs/architecture.md"): (
        "least-authorized registered route",
        "does not consult",
        "credential resolution",
    ),
    Path("docs/configuration.md"): (
        "Authenticated GitHub Cloud capabilities",
        "never resolves or sends that token",
        "Atlassian account email and",
        "MASTER_AGENT_BITBUCKET_EMAIL",
        "bitbucket.public_repository.list",
        "constructs an anonymous",
        "sends ambient Bitbucket credentials",
    ),
    Path("docs/cli-reference.md"): (
        "github-repositories --username USERNAME",
        "ignores ambient GitHub tokens",
        "bitbucket-repositories --workspace WORKSPACE",
        "constructs an anonymous Bitbucket Cloud connector",
    ),
    Path("docs/copilot-custom-agent.md"): (
        "github-repositories --username USERNAME",
        "credential-free typed route",
        "bitbucket-repositories --workspace WORKSPACE",
        "ignores ambient Bitbucket credentials",
    ),
    Path("docs/deployment-runbook.md"): (
        "do not provision",
        "github-repositories --username USERNAME",
        "bitbucket-repositories --workspace WORKSPACE",
        "ignore ambient Bitbucket credentials",
    ),
    Path("docs/implementation-roadmap.md"): (
        "github.public_repository.list",
        "bitbucket.public_repository.list",
        "operate anonymously",
    ),
    Path("docs/github-connector-quickstart.md"): (
        "Public repositories need no credential",
        "loads or sends a token",
    ),
    Path("docs/integration-matrix.md"): (
        "Cloud anonymous public reads",
        "anonymous route omits credentials",
        "public-workspace repository lists",
    ),
    Path("docs/live-connectors.md"): (
        "github.public_repository.list",
        "does not resolve or forward an ambient credential",
        "bitbucket.public_repository.list",
        "Bitbucket Cloud workspace-repositories endpoint",
    ),
    Path("docs/operations.md"): (
        "authentication class is `none`",
        "must not resolve one",
    ),
    Path("docs/phase-2-read-only.md"): (
        "explicitly cataloged for anonymous access",
        "must not resolve or forward ambient credentials",
        "bitbucket.public_repository.list",
    ),
    Path("docs/phase-2c-authentication.md"): (
        "cataloged for anonymous public access",
        "does not acquire, resolve, or",
        "bitbucket.public_repository.list",
    ),
    Path("docs/release-validation.md"): (
        "Typed anonymous public-data capabilities require no credential activation",
        "rejects the stale blanket claims",
    ),
    Path("docs/threat-model.md"): (
        "typed anonymous public-data",
        "never resolves or forwards an ambient credential",
    ),
}

_PUBLIC_READ_FORBIDDEN_CLAIMS = {
    Path("README.md"): (
        "Organization-approved HTTPS API endpoints and credentials for live use.",
    ),
    Path("docs/configuration.md"): (
        "GitHub Cloud uses `MASTER_AGENT_GITHUB_TOKEN` as a bearer token.",
    ),
    Path("docs/copilot-custom-agent.md"): (
        "Bitbucket, GitHub, Microsoft identity, SharePoint, Outlook, Teams, or OneNote",
    ),
    Path("docs/deployment-runbook.md"): (
        "## 4. Register applications and credentials",
    ),
    Path("docs/live-connectors.md"): (
        "GitHub Cloud connector constructs authenticated-user repository-list",
    ),
    Path("docs/phase-2-read-only.md"): (
        "Available only for explicitly unauthenticated test or internal endpoints.",
    ),
    Path("docs/phase-2c-authentication.md"): (
        "Before live use, administrators must review:",
    ),
    Path("docs/release-validation.md"): ("Before live use, administrators must",),
}

_RETENTION_PRUNE_DOCUMENT_REQUIREMENTS = {
    Path("CHANGELOG.md"): (
        "Enable descriptor-safe expiration deletion",
        "preview, apply, and orphan repair remain capability-gated",
    ),
    Path("docs/architecture.md"): (
        "Expiration preview and explicit POSIX deletion share",
        "every discovered evidence-parent lock",
        "content-free transaction before either public name",
    ),
    Path("docs/cli-reference.md"): (
        "## Evidence expiration maintenance",
        "every discovered evidence-parent retention lock",
        "All `evidence-prune` execution remains unavailable on Windows",
        "Windows preview and apply are capability-gated",
    ),
    Path("docs/configuration.md"): (
        "## Retention and expiry",
        "Changing `retention.toml` later does not",
        "root and discovered evidence-parent locks",
    ),
    Path("docs/implementation-roadmap.md"): (
        "Complete on POSIX; native Windows retention preview/apply/repair gated",
        "native Windows secure filesystem/ACL, locking, atomic state and retention",
    ),
    Path("docs/operations.md"): (
        "repeat the apply command under the same root",
        "bounded and uses the same evidence-parent locks",
        "All `evidence-prune` execution remains unavailable on Windows",
        "Windows retention preview, apply, and orphan repair remain unavailable",
    ),
    Path("docs/phase-2b-communication-context.md"): (
        "same bounded descriptor-relative validation plan",
        "discovered evidence-parent retention lock",
    ),
    Path("docs/release-validation.md"): (
        "POSIX retained-evidence expiration tests prove",
        "apply, and orphan repair remain capability-gated",
    ),
    Path("docs/threat-model.md"): (
        "every discovered evidence-parent",
        "broad, path-based, or unvalidated recursive evidence deletion",
        "common Windows import and configuration-diagnostics boundary is not",
        "orphan repair remain unavailable",
    ),
}

_RETENTION_PRUNE_FORBIDDEN_CURRENT_CLAIMS = {
    Path("docs/architecture.md"): ("expiry deletion remains preview-only",),
    Path("docs/cli-reference.md"): (
        "`--apply` is disabled before traversal or deletion",
    ),
    Path("docs/implementation-roadmap.md"): ("destructive retention pruning",),
    Path("docs/operations.md"): ("Evidence expiry deletion remains preview-only",),
    Path("docs/phase-2b-communication-context.md"): (
        "Expiry deletion remains preview-only",
        "`evidence-prune --apply` still fails before traversal",
    ),
    Path("docs/threat-model.md"): (
        "expiry deletion remains preview-only",
        "expiry deletion is preview-only",
        "destructive recursive evidence pruning",
    ),
}

_CAPSULE_DOCUMENT_REQUIREMENTS = {
    Path(".ai/MASTER_AGENT.md"): (
        "newly generated capability code as quarantined data",
        "Never let generated code sign, review, publish, enable, route, approve",
        "docs/capability-capsules.md",
    ),
    _AUTONOMY_CONTRACT_PATH: (
        "Implementation does not make generated code trusted",
        "docs/capability-capsules.md",
    ),
    Path("README.md"): (
        "Capability capsule promotion",
        "Linux bubblewrap",
        "docs/capability-capsules.md",
    ),
    Path("docs/capability-capsules.md"): (
        "dependency-free, deterministic",
        "## Isolation boundary",
        "## Credential broker boundary",
        "## Supply-chain admission",
        "## Production boundary",
        "shipped deployment remains\nfail closed",
    ),
    Path("docs/architecture.md"): (
        "### Capability capsule promotion",
        "Linux bubblewrap",
        "capability-capsules.md",
    ),
    Path("docs/threat-model.md"): (
        "Generated capability substitution or self-promotion",
        "AST-restricted language in Linux bubblewrap",
        "capability-capsules.md",
    ),
    Path("docs/release-validation.md"): (
        "capability-capsule acceptance flow",
        "dependency-license admission policy",
        "CycloneDX 1.5 SBOM",
    ),
}

_ADVISORY_DOCUMENT_REQUIREMENTS = {
    Path("AGENTS.md"): (
        "Direct GitHub-host advisory sub-agent invocation is disabled",
        "repository-owned advisory integration harness",
        "complete the same work directly",
        "`--goal-id`",
        "`--route ROUTE_ID`",
        "`--path`",
    ),
    Path(".ai/MASTER_AGENT.md"): (
        "Direct GitHub-host advisory sub-agent invocation is disabled",
        "repository-owned advisory integration harness",
        "complete the same work directly",
        "`--goal-id`",
        "`--route ROUTE_ID`",
        "untracked-byte",
    ),
    _AUTONOMY_CONTRACT_PATH: (
        "## Bounded advisory delegation",
        "Direct GitHub-host invocation is disabled",
        "at most three research attempts and one plan review",
        "complete the same work directly",
        "`scripts/advisory_subagent.py`",
        "`--route ROUTE_ID`",
    ),
    _COPILOT_AGENT_PATH: (
        "## Advisory boundary",
        "Direct GitHub-host advisory invocation is disabled",
        "repository-owned advisory integration harness",
        "complete the same work directly",
        "optional current Copilot SDK adapter",
        "`--route ROUTE_ID`",
    ),
    Path("README.md"): (
        "checked-in advisory profiles now define a fail-closed contract",
        "repository-owned advisory integration harness",
        "completes the same research or review directly",
        "authenticated cross-process goal budget",
        "`--route ROUTE_ID`",
    ),
    Path("CHANGELOG.md"): (
        "Harden advisory sub-agent boundaries end to end",
        "researcher no longer has generic execution",
        "HMAC-authenticated cross-process goal",
    ),
    Path("docs/advisory-subagents.md"): (
        "## Repository-owned integration harness",
        "## Hermetic end-to-end tests",
        "## Goal budget",
        "## Technical route scope",
        "`--route ROUTE_ID`",
        "No live Copilot canary is bundled",
        "The deterministic runtime remains the only path",
    ),
    Path("docs/architecture.md"): (
        "no direct GitHub-host child invocation",
        "`src/master_agent/advisory.py` loads those exact profiles",
        "`src/master_agent/advisory_budget.py`",
        "explicit parent fallback",
    ),
    Path("docs/copilot-custom-agent.md"): (
        "direct GitHub-host invocation is disabled",
        "repository-owned integration harness",
        "completes the same work directly",
        "`advisory_subagent.py`",
        "`--route ROUTE_ID`",
    ),
    Path("docs/release-validation.md"): (
        "Direct child user/model invocation",
        "profile-derived dispatch",
        "`--route ROUTE_ID`",
        "untracked-file mutation",
        "no filesystem",
    ),
    Path("docs/threat-model.md"): (
        "direct GitHub-host advisory invocation is disabled",
        "profile-derived dispatcher denies execute",
        "HMAC-authenticated state",
        "bounded untracked file",
        "falls back to the parent without changing filesystem",
    ),
}


@dataclass(frozen=True, slots=True)
class ValidationReport:
    """Result of a release-tree validation.

    Parameters
    ----------
    checks
        Human-readable successful checks.
    errors
        Human-readable validation failures.
    """

    checks: tuple[str, ...]
    errors: tuple[str, ...]

    @property
    def valid(self) -> bool:
        """Return whether the source tree passed every check."""

        return not self.errors

    def to_dict(self) -> dict[str, object]:
        """Serialize the report to JSON-compatible data."""

        return {
            "schema": "master-agent/release-validation@1",
            "valid": self.valid,
            "checks": list(self.checks),
            "errors": list(self.errors),
        }


def validate_project(root: Path) -> ValidationReport:
    """Validate release metadata, safe defaults, links, and file hygiene.

    Parameters
    ----------
    root
        Project root containing ``pyproject.toml``.

    Returns
    -------
    ValidationReport
        Completed validation report.
    """

    root = root.resolve()
    checks: list[str] = []
    errors: list[str] = []

    _validate_versions(root, checks, errors)
    _validate_supply_chain(root, checks, errors)
    _validate_packaged_defaults(root, checks, errors)
    _validate_capabilities(root, checks, errors)
    _validate_copilot_agent(root, checks, errors)
    _validate_advisory_agents(root, checks, errors)
    _validate_advisory_contract(root, checks, errors)
    _validate_semantic_router(root, checks, errors)
    _validate_first_run_contract(root, checks, errors)
    _validate_public_read_contract(root, checks, errors)
    _validate_capsule_contract(root, checks, errors)
    _validate_retention_prune_contract(root, checks, errors)
    _validate_markdown_links(root, checks, errors)
    _validate_documentation(root, checks, errors)
    _validate_demo(root, checks, errors)
    _validate_file_hygiene(root, checks, errors)

    return ValidationReport(tuple(checks), tuple(errors))


def validate_archive(path: Path) -> ValidationReport:
    """Validate one built wheel or source archive without extracting it."""

    checks: list[str] = []
    errors: list[str] = []
    if not path.is_file():
        return ValidationReport((), (f"release archive not found: {path}",))
    names: list[str] = []
    try:
        if path.suffix == ".whl":
            with zipfile.ZipFile(path) as archive:
                bad_member = archive.testzip()
                if bad_member is not None:
                    errors.append(f"wheel integrity check failed: {bad_member}")
                for item in archive.infolist():
                    if item.is_dir():
                        continue
                    names.append(item.filename)
                    archive_mode = (item.external_attr >> 16) & 0xFFFF
                    if stat.S_ISLNK(archive_mode):
                        errors.append(
                            f"release archive contains link entry: {item.filename}"
                        )
                    _validate_archive_runtime_mode(
                        item.filename,
                        archive_mode,
                        is_regular=stat.S_ISREG(archive_mode),
                        errors=errors,
                    )
                    with archive.open(item) as handle:
                        _consume_stream(handle)
        elif path.name.endswith((".tar.gz", ".tar.bz2", ".tar.xz", ".tgz")):
            with tarfile.open(path, mode="r:*") as archive:
                for item in archive.getmembers():
                    if item.issym() or item.islnk():
                        errors.append(
                            f"release archive contains link entry: {item.name}"
                        )
                        continue
                    if not item.isfile():
                        continue
                    names.append(item.name)
                    _validate_archive_runtime_mode(
                        item.name,
                        item.mode,
                        is_regular=True,
                        errors=errors,
                    )
                    handle = archive.extractfile(item)
                    if handle is not None:
                        with handle:
                            _consume_stream(handle)
        else:
            errors.append(f"unsupported release archive type: {path.name}")
    except (OSError, tarfile.TarError, zipfile.BadZipFile) as error:
        errors.append(f"release archive could not be read: {path}: {error}")

    for name in names:
        _validate_archive_member(name, errors)
    if path.suffix == ".whl":
        required_members = (
            "master_agent/__init__.py",
            "master_agent/capsule_worker.py",
            "master_agent/platform_runtime/posix/capsule_worker.py",
            "master_agent/defaults/capabilities.toml",
            "master_agent/defaults/dependency-licenses.toml",
        )
        for required in required_members:
            if required not in names:
                errors.append(f"release archive is missing required file: {required}")
        wheel_parts = path.name.removesuffix(".whl").split("-")
        metadata_member = (
            f"master_agent-{wheel_parts[1]}.dist-info/METADATA"
            if len(wheel_parts) >= 5 and wheel_parts[0] == "master_agent"
            else None
        )
        if metadata_member is None or metadata_member not in names:
            errors.append(
                "release archive is missing required file: "
                "master_agent-<version>.dist-info/METADATA"
            )
    else:
        required_paths = (
            "/.ai/MASTER_AGENT.md",
            "/.ai/FIRST_RUN.md",
            "/.ai/AUTONOMY.md",
            "/.ai/semantic-router.toml",
            "/.github/agents/MasterAgent.agent.md",
            "/.github/agents/MasterAgent-Read-Researcher.agent.md",
            "/.github/agents/MasterAgent-Plan-Reviewer.agent.md",
            "/.github/workflows/ci.yml",
            "/.github/workflows/confluence-sandbox.yml",
            "/.github/workflows/github-actions-live-integration.yml",
            "/.github/workflows/live-connector-integration.yml",
            "/.env.example",
            "/LICENSE",
            "/setup.py",
            "/THIRD_PARTY_NOTICES.md",
            "/config/capabilities.toml",
            "/config/dependency-licenses.toml",
            "/requirements-runtime.lock",
            "/sbom.cdx.json",
            "/supply-chain/runtime-dependencies.toml",
            "/docs/capability-capsules.md",
            "/docs/semantic-index.md",
            "/docs/semantic-router-metrics.md",
            "/scripts/bootstrap_agent.py",
            "/scripts/generate_sbom.py",
            "/scripts/semantic_router.py",
            "/scripts/validate_release.py",
            "/tests/test_capability_capsules.py",
            "/tests/test_capsule_broker_and_routing.py",
            "/tests/test_release_metadata.py",
            "/tests/test_live_connector_workflow.py",
            "/tests/test_semantic_router.py",
            "/tests/test_advisory_integration.py",
            "/tests/fixtures/advisory/repository_prompt_injection.txt",
            "/tests/fixtures/advisory/provider_prompt_injection.txt",
            "/specs/current/security/MA-ADVISORY-001.md",
            "/specs/current/development/MA-ROUTER-001.md",
            "/src/master_agent/__init__.py",
            "/src/master_agent/advisory.py",
            "/src/master_agent/capsule_worker.py",
            "/src/master_agent/platform_runtime/posix/capsule_worker.py",
        )
        source_members = tuple(PurePosixPath(name) for name in names)
        source_roots = {
            member.parts[0] for member in source_members if len(member.parts) >= 2
        }
        common_root = (
            next(iter(source_roots))
            if source_members
            and len(source_roots) == 1
            and all(len(member.parts) >= 2 for member in source_members)
            else None
        )
        if common_root is None:
            errors.append(
                "source archive must contain exactly one top-level root directory"
            )
        for required in required_paths:
            required_parts = PurePosixPath(required.removeprefix("/")).parts
            expected_parts = (
                (common_root, *required_parts) if common_root is not None else None
            )
            if expected_parts is None or not any(
                member.parts == expected_parts for member in source_members
            ):
                errors.append(f"release archive is missing required file: {required}")
    if not errors:
        checks.append(
            f"validated release archive {path.name} ({len(names)} files, no links)"
        )
    return ValidationReport(tuple(checks), tuple(errors))


def _validate_semantic_router(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Validate exact semantic ownership, topology, and generated output."""

    try:
        manifest = load_semantic_manifest(root)
        semantic_errors = validate_semantic_repository(root, manifest)
    except (SemanticManifestError, OSError) as error:
        errors.append(f"semantic router could not be validated: {error}")
        return
    if semantic_errors:
        errors.extend(f"semantic router: {error}" for error in semantic_errors)
        return
    checks.append(
        "semantic router covers "
        f"{len(manifest.routes)} routes and "
        f"{len(manifest.routing_cases)} routing fixtures"
    )


def sha256_file(path: Path) -> str:
    """Return a file's SHA-256 digest.

    Parameters
    ----------
    path
        File to hash.

    Returns
    -------
    str
        Lowercase hexadecimal digest.
    """

    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _validate_versions(root: Path, checks: list[str], errors: list[str]) -> None:
    pyproject = tomllib.loads((root / "pyproject.toml").read_text(encoding="utf-8"))
    project_version = str(pyproject["project"]["version"])
    init_text = (root / "src/master_agent/__init__.py").read_text(encoding="utf-8")
    match = re.search(r'__version__\s*=\s*"([^"]+)"', init_text)
    module_version = match.group(1) if match else ""
    http_text = (root / "src/master_agent/http.py").read_text(encoding="utf-8")
    if project_version != module_version:
        errors.append(
            f"version mismatch: pyproject={project_version}, module={module_version}"
        )
    elif f"master-agent/{project_version}" not in http_text:
        errors.append("HTTP user agent does not match the project version")
    else:
        checks.append(f"version metadata is consistent at {project_version}")


def _validate_supply_chain(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    error_count = len(errors)
    required = (
        "LICENSE",
        "THIRD_PARTY_NOTICES.md",
        "requirements-runtime.lock",
        "sbom.cdx.json",
        "config/dependency-licenses.toml",
        "supply-chain/runtime-dependencies.toml",
        "scripts/generate_sbom.py",
    )
    missing = [name for name in required if not (root / name).is_file()]
    if missing:
        errors.append("supply-chain baseline files are missing: " + ", ".join(missing))
        return
    pyproject = tomllib.loads((root / "pyproject.toml").read_text(encoding="utf-8"))
    project_metadata = pyproject["project"]
    if project_metadata.get(
        "license"
    ) != "LicenseRef-MasterAgent-Proprietary" or project_metadata.get(
        "license-files"
    ) != ["LICENSE"]:
        errors.append("pyproject license does not bind the repository LICENSE")
    policy_document = tomllib.loads(
        (root / "config/dependency-licenses.toml").read_text(encoding="utf-8")
    )
    policy = policy_document.get("policy")
    if not isinstance(policy, dict):
        errors.append("dependency-license policy is missing or malformed")
        return
    allowed = policy.get("allowed_spdx")
    denied = policy.get("denied_spdx")
    if (
        not isinstance(allowed, list)
        or not all(isinstance(value, str) and value for value in allowed)
        or not isinstance(denied, list)
        or not all(isinstance(value, str) and value for value in denied)
        or set(allowed) & set(denied)
        or not isinstance(policy.get("deny_unknown"), bool)
        or not isinstance(policy.get("require_notices"), bool)
    ):
        errors.append("dependency-license policy is missing or malformed")
        return
    inventory = tomllib.loads(
        (root / "supply-chain/runtime-dependencies.toml").read_text(encoding="utf-8")
    )
    project = inventory.get("project")
    components = inventory.get("components", [])
    if (
        not isinstance(project, dict)
        or not isinstance(components, list)
        or not components
    ):
        errors.append("runtime dependency inventory is empty or malformed")
        return
    required_component_fields = {
        "name",
        "version",
        "license",
        "purl",
        "homepage",
        "notice",
        "dependencies",
    }
    if any(
        not isinstance(item, dict)
        or not required_component_fields.issubset(item)
        or any(
            not isinstance(item.get(field), str) or not str(item[field]).strip()
            for field in required_component_fields - {"dependencies"}
        )
        or not isinstance(item.get("dependencies"), list)
        or not all(
            isinstance(value, str) and value for value in item.get("dependencies", [])
        )
        for item in components
    ):
        errors.append("runtime dependency inventory is empty or malformed")
        return
    expected = {
        str(item["name"]).replace("_", "-").casefold(): (
            str(item["version"]),
            str(item["license"]),
        )
        for item in components
    }
    if len(expected) != len(components):
        errors.append("runtime dependency inventory contains duplicate names")
    observed_licenses = {
        str(project.get("license", "")),
        *(str(item["license"]) for item in components),
    }
    blocked_licenses = sorted(observed_licenses & set(denied))
    unknown_licenses = sorted(observed_licenses - set(allowed))
    if blocked_licenses:
        errors.append(
            "runtime dependency license is denied: " + ", ".join(blocked_licenses)
        )
    if bool(policy["deny_unknown"]) and unknown_licenses:
        errors.append(
            "runtime dependency license is unknown: " + ", ".join(unknown_licenses)
        )
    if bool(policy["require_notices"]) and any(
        not str(item["notice"]).strip() for item in components
    ):
        errors.append("runtime dependency notice is missing")

    direct = project.get("dependencies")
    optional_extra = project.get("optional_extra")
    optional_dependencies = project.get("optional_dependencies")
    dependency_scopes_are_valid = (
        isinstance(direct, list)
        and all(isinstance(value, str) and value for value in direct)
        and isinstance(optional_extra, str)
        and bool(optional_extra)
        and isinstance(optional_dependencies, list)
        and bool(optional_dependencies)
        and all(isinstance(value, str) and value for value in optional_dependencies)
    )
    if not dependency_scopes_are_valid:
        errors.append("project dependency scopes are malformed")
    else:
        by_name = {
            str(item["name"]).replace("_", "-").casefold(): item for item in components
        }
        pending = [
            value.replace("_", "-").casefold()
            for value in (*direct, *optional_dependencies)
        ]
        reached: set[str] = set()
        while pending:
            name = pending.pop()
            if name in reached:
                continue
            component = by_name.get(name)
            if component is None:
                errors.append(f"runtime dependency closure is missing: {name}")
                break
            reached.add(name)
            pending.extend(
                value.replace("_", "-").casefold()
                for value in component["dependencies"]
            )
        if reached != set(by_name):
            errors.append("runtime dependency inventory is not one complete closure")

    declared: dict[str, str] = {}
    for requirement in pyproject["project"].get("dependencies", []):
        match = re.fullmatch(r"([A-Za-z0-9_.-]+)==([A-Za-z0-9_.+-]+)", requirement)
        if match is None:
            errors.append(f"runtime project dependency is not exact: {requirement}")
            continue
        declared[match.group(1).replace("_", "-").casefold()] = match.group(2)
    if dependency_scopes_are_valid:
        expected_core = {
            name.replace("_", "-").casefold(): expected.get(
                name.replace("_", "-").casefold(), ("", "")
            )[0]
            for name in direct
        }
        if declared != expected_core:
            errors.append("pyproject core dependencies differ from the inventory")

        optional_metadata = project_metadata.get("optional-dependencies")
        optional_declared: dict[str, str] = {}
        optional_requirements = (
            optional_metadata.get(optional_extra)
            if isinstance(optional_metadata, dict)
            else None
        )
        if not isinstance(optional_requirements, list) or not all(
            isinstance(requirement, str) for requirement in optional_requirements
        ):
            errors.append("pyproject optional draft dependencies are malformed")
        else:
            for requirement in optional_requirements:
                match = re.fullmatch(
                    r"([A-Za-z0-9_.-]+)==([A-Za-z0-9_.+-]+)", requirement
                )
                if match is None:
                    errors.append(
                        f"optional draft dependency is not exact: {requirement}"
                    )
                    continue
                optional_declared[match.group(1).replace("_", "-").casefold()] = (
                    match.group(2)
                )
            expected_optional = {
                name.replace("_", "-").casefold(): expected.get(
                    name.replace("_", "-").casefold(), ("", "")
                )[0]
                for name in optional_dependencies
            }
            if optional_declared != expected_optional:
                errors.append(
                    "pyproject optional draft dependencies differ from the inventory"
                )
    lock_entries: dict[str, str] = {}
    for line in (
        (root / "requirements-runtime.lock").read_text(encoding="utf-8").splitlines()
    ):
        if not line or line.startswith("#"):
            continue
        match = re.fullmatch(r"([A-Za-z0-9_.-]+)==([A-Za-z0-9_.+-]+)", line)
        if match is None:
            errors.append(f"runtime lock entry is not exact: {line}")
            continue
        lock_entries[match.group(1).replace("_", "-").casefold()] = match.group(2)
    if lock_entries != {name: value[0] for name, value in expected.items()}:
        errors.append("runtime lock differs from the reviewed dependency inventory")
    sbom = json.loads((root / "sbom.cdx.json").read_text(encoding="utf-8"))
    sbom_components = sbom.get("components", []) if isinstance(sbom, dict) else []
    observed = {
        str(item.get("name", "")).replace("_", "-").casefold(): (
            str(item.get("version", "")),
            str(item.get("licenses", [{}])[0].get("license", {}).get("id", "")),
        )
        for item in sbom_components
        if isinstance(item, dict)
        and isinstance(item.get("licenses"), list)
        and item.get("licenses")
        and isinstance(item.get("licenses", [None])[0], dict)
        and isinstance(item.get("licenses", [{}])[0].get("license"), dict)
    }
    metadata = sbom.get("metadata") if isinstance(sbom, dict) else None
    metadata_component = (
        metadata.get("component") if isinstance(metadata, dict) else None
    )
    properties = (
        metadata_component.get("properties")
        if isinstance(metadata_component, dict)
        else None
    )
    marks_optional_drafts = isinstance(properties, list) and any(
        isinstance(property_, dict)
        and property_.get("name") == "master-agent:optional-extra"
        and property_.get("value") == optional_extra
        for property_ in properties
    )
    if (
        not isinstance(sbom, dict)
        or sbom.get("bomFormat") != "CycloneDX"
        or str(sbom.get("specVersion")) != "1.5"
        or observed != expected
        or not dependency_scopes_are_valid
        or not marks_optional_drafts
    ):
        errors.append("CycloneDX SBOM differs from the optional draft dependency lock")
    notices = (root / "THIRD_PARTY_NOTICES.md").read_text(encoding="utf-8")
    if any(
        str(item.get("name", "")) not in notices
        or str(item.get("license", "")) not in notices
        for item in components
        if isinstance(item, dict)
    ):
        errors.append("third-party notices omit a runtime component or license")
    if len(errors) == error_count:
        checks.append(
            f"license, policy, exact lock, CycloneDX SBOM, and notices cover "
            f"{len(expected)} optional draft components"
        )


def _validate_packaged_defaults(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    config_dir = root / "config"
    defaults_dir = root / "src/master_agent/defaults"
    expected = sorted(path.name for path in config_dir.glob("*.toml"))
    for name in expected:
        source = config_dir / name
        packaged = defaults_dir / name
        if not packaged.exists():
            errors.append(f"packaged default is missing: {name}")
            continue
        if source.read_bytes() != packaged.read_bytes():
            errors.append(f"packaged default differs from repository config: {name}")
    if not errors:
        checks.append(f"{len(expected)} packaged TOML defaults match repository config")

    integrations = tomllib.loads((defaults_dir / "integrations.toml").read_text())
    connectors = integrations.get("connectors", {})
    read_gates = {("microsoft", "onenote_read_enabled")}
    for name, connector in connectors.items():
        if not bool(connector.get("enabled", False)):
            errors.append(f"packaged connector is unavailable: {name}")
        for key, value in connector.items():
            if (
                key.endswith("_enabled")
                and bool(value)
                and (name, key) not in read_gates
                and key != "enabled"
            ):
                errors.append(f"packaged provider gate is enabled: {name}.{key}")
    if connectors and all(
        bool(value.get("enabled", False)) for value in connectors.values()
    ):
        checks.append(
            "all packaged read connectors are available and mutation gates are disabled"
        )

    recurring = tomllib.loads((defaults_dir / "recurring.toml").read_text())
    workflows = recurring.get("workflows", {})
    enabled = [name for name, item in workflows.items() if item.get("enabled")]
    if enabled:
        errors.append(f"packaged recurring workflows are enabled: {enabled}")
    else:
        checks.append("all packaged recurring workflows are disabled")


def _validate_capabilities(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    raw = tomllib.loads((root / "config/capabilities.toml").read_text())
    capabilities = raw.get("capabilities", {})
    if len(capabilities) != 82:
        errors.append(f"expected 82 v1 capabilities, found {len(capabilities)}")
    else:
        checks.append("capability catalog contains 82 typed capabilities")
    merge = capabilities.get("bitbucket.pull_request.merge", {})
    if merge.get("enabled") is not False:
        errors.append("Bitbucket pull-request merge must remain disabled")
    else:
        checks.append("high-impact pull-request merge remains disabled")


def _validate_copilot_agent(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Require the repository-scoped Copilot entry point to remain fail-closed."""

    path = root / _COPILOT_AGENT_PATH
    try:
        text = path.read_text(encoding="utf-8")
    except (OSError, UnicodeError) as error:
        errors.append(f"Copilot custom agent is missing or unreadable: {path}: {error}")
        return

    metadata, body, frontmatter_errors = _parse_agent_frontmatter(text)
    errors.extend(
        f"Copilot custom agent frontmatter: {error}" for error in frontmatter_errors
    )
    if frontmatter_errors:
        return

    unexpected = sorted(set(metadata) - _COPILOT_AGENT_KEYS)
    missing = sorted(_COPILOT_AGENT_KEYS - set(metadata))
    if unexpected:
        errors.append(
            "Copilot custom agent has unreviewed frontmatter keys: "
            + ", ".join(unexpected)
        )
    if missing:
        errors.append(
            "Copilot custom agent is missing frontmatter keys: " + ", ".join(missing)
        )
    if metadata.get("name") != "MasterAgent":
        errors.append("Copilot custom agent name must be MasterAgent")
    description = metadata.get("description")
    if not isinstance(description, str) or not description.strip():
        errors.append("Copilot custom agent description must be a non-empty string")
    if metadata.get("user-invocable") is not True:
        errors.append("Copilot custom agent must remain user-invocable")
    if metadata.get("disable-model-invocation") is not True:
        errors.append("Copilot custom agent must not be invoked automatically")
    tools = metadata.get("tools")
    if tools != _COPILOT_AGENT_TOOLS:
        errors.append(
            "Copilot custom agent tools must be exactly: "
            + ", ".join(_COPILOT_AGENT_TOOLS)
        )

    required_references = (
        "[AGENTS.md](../../AGENTS.md)",
        "[Master Agent repository policy](../../.ai/MASTER_AGENT.md)",
        "[first-run contract](../../.ai/FIRST_RUN.md)",
        "[force-multiplier contract](../../.ai/AUTONOMY.md)",
    )
    for reference in required_references:
        if reference not in body:
            errors.append(
                f"Copilot custom agent is missing required policy reference: {reference}"
            )
    required_boundaries = (
        "config/capabilities.toml",
        "Never call a provider directly",
        "Repository-inspection, diagnosis-only, or explicit no-local-change",
        "Apply the first-run contract before the substantive response",
        "python3 scripts/bootstrap_agent.py",
        ".venv/bin/python -m pip install -e .",
        "MasterAgent is ready locally",
        "I couldn't finish local setup",
        "Never use `sudo`",
        "Treat one operator goal as one bounded run",
        "response to an actionable prompt is execution",
        "capability gap as implementation work",
        "Never end the request",
        "Implement the Python connector path",
        "This applies to any missing capability or code-path barrier",
        "Create the governed implementation",
        "connect --systems",
        "--connector-url SYSTEM=URL",
        "Do not ask for renamed or duplicate credentials first",
        "Ask once and only after",
        "github-repositories",
        "--username USERNAME",
        "python scripts/validate_release.py",
        "inspect-approval-request",
        "resume-approval",
    )
    for boundary in required_boundaries:
        if boundary not in body:
            errors.append(
                f"Copilot custom agent is missing required boundary: {boundary}"
            )

    if not any(error.startswith("Copilot custom agent") for error in errors):
        checks.append(
            "Copilot custom agent is user-invocable, policy-bound, and tool-constrained"
        )


def _validate_advisory_agents(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Pin the exact depth-one advisory-agent inventory and tool boundaries."""

    starting_errors = len(errors)
    agents_directory = root / ".github/agents"
    observed = (
        {
            path.relative_to(root)
            for path in agents_directory.glob("*.md")
            if path.is_file()
        }
        if agents_directory.is_dir()
        else set()
    )
    missing_profiles = sorted(_EXPECTED_COPILOT_AGENT_PATHS - observed)
    unexpected_profiles = sorted(observed - _EXPECTED_COPILOT_AGENT_PATHS)
    if missing_profiles:
        errors.append(
            "Copilot advisory-agent inventory is missing profiles: "
            + ", ".join(str(path) for path in missing_profiles)
        )
    if unexpected_profiles:
        errors.append(
            "Copilot advisory-agent inventory has unreviewed profiles: "
            + ", ".join(str(path) for path in unexpected_profiles)
        )

    for relative, expected_name in _ADVISORY_AGENT_NAMES.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"Copilot advisory agent is missing or unreadable: {relative}: {error}"
            )
            continue

        metadata, body, frontmatter_errors = _parse_agent_frontmatter(text)
        errors.extend(
            f"Copilot advisory agent {relative} frontmatter: {error}"
            for error in frontmatter_errors
        )
        if frontmatter_errors:
            continue

        unexpected_keys = sorted(set(metadata) - _COPILOT_AGENT_KEYS)
        missing_keys = sorted(_COPILOT_AGENT_KEYS - set(metadata))
        if unexpected_keys:
            errors.append(
                f"Copilot advisory agent {relative} has unreviewed frontmatter "
                "keys: " + ", ".join(unexpected_keys)
            )
        if missing_keys:
            errors.append(
                f"Copilot advisory agent {relative} is missing frontmatter keys: "
                + ", ".join(missing_keys)
            )
        if metadata.get("name") != expected_name:
            errors.append(
                f"Copilot advisory agent {relative} must be named {expected_name}"
            )
        description = metadata.get("description")
        if not isinstance(description, str) or not description.strip():
            errors.append(
                f"Copilot advisory agent {relative} description must be non-empty"
            )
        if metadata.get("user-invocable") is not False:
            errors.append(
                f"Copilot advisory agent {relative} direct user invocation must remain disabled"
            )
        if metadata.get("disable-model-invocation") is not True:
            errors.append(
                f"Copilot advisory agent {relative} direct model invocation must remain disabled"
            )
        tools = metadata.get("tools")
        expected_tools = _ADVISORY_AGENT_TOOLS[relative]
        if tools != expected_tools:
            errors.append(
                f"Copilot advisory agent {relative} tools must be exactly: "
                + ", ".join(expected_tools)
            )
        for boundary in _ADVISORY_AGENT_BOUNDARIES[relative]:
            if boundary not in body:
                errors.append(
                    f"Copilot advisory agent {relative} is missing required "
                    f"boundary: {boundary}"
                )

    if len(errors) == starting_errors:
        checks.append(
            "Copilot advisory profiles are non-invocable, read/search-only, and "
            "fail-closed"
        )


def _validate_advisory_contract(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Keep advisory delegation policy consistent across durable guidance."""

    starting_errors = len(errors)
    for relative, requirements in _ADVISORY_DOCUMENT_REQUIREMENTS.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"advisory sub-agent contract document is unreadable: "
                f"{relative}: {error}"
            )
            continue
        for requirement in requirements:
            if requirement not in text:
                errors.append(
                    "advisory sub-agent contract document is inconsistent: "
                    f"{relative} is missing {requirement!r}"
                )

    if len(errors) == starting_errors:
        checks.append(
            "advisory sub-agent safety guidance is consistent across "
            f"{len(_ADVISORY_DOCUMENT_REQUIREMENTS)} policy and documentation files"
        )


def _validate_first_run_contract(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Keep every first-run instruction and onboarding guide synchronized."""

    starting_errors = len(errors)
    for relative, requirements in _FIRST_RUN_DOCUMENT_REQUIREMENTS.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"first-run contract document is unreadable: {relative}: {error}"
            )
            continue
        for requirement in requirements:
            if requirement not in text:
                errors.append(
                    "first-run contract document is inconsistent: "
                    f"{relative} is missing {requirement!r}"
                )

    script = root / "scripts/bootstrap_agent.py"
    if not script.is_file():
        errors.append(
            "first-run bootstrap script is missing: scripts/bootstrap_agent.py"
        )

    if len(errors) == starting_errors:
        checks.append(
            "first-run and force-multiplier contracts are consistent across 9 instruction and onboarding files"
        )


def _validate_public_read_contract(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Keep anonymous public reads distinct from authenticated provider access."""

    starting_errors = len(errors)
    for relative, requirements in _PUBLIC_READ_DOCUMENT_REQUIREMENTS.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"public-read contract document is unreadable: {relative}: {error}"
            )
            continue
        for requirement in requirements:
            if requirement not in text:
                errors.append(
                    "public-read contract document is inconsistent: "
                    f"{relative} is missing {requirement!r}"
                )
        for forbidden in _PUBLIC_READ_FORBIDDEN_CLAIMS.get(relative, ()):
            if forbidden in text:
                errors.append(
                    "public-read contract contains a blanket credential requirement: "
                    f"{relative} contains {forbidden!r}"
                )

    if len(errors) == starting_errors:
        checks.append(
            "anonymous public-read guidance is consistent across "
            f"{len(_PUBLIC_READ_DOCUMENT_REQUIREMENTS)} policy and documentation files"
        )


def _validate_capsule_contract(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Keep generated-code safety boundaries consistent across public guidance."""

    starting_errors = len(errors)
    for relative, requirements in _CAPSULE_DOCUMENT_REQUIREMENTS.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"capability-capsule contract document is unreadable: {relative}: {error}"
            )
            continue
        for requirement in requirements:
            if requirement not in text:
                errors.append(
                    "capability-capsule contract document is inconsistent: "
                    f"{relative} is missing {requirement!r}"
                )

    if len(errors) == starting_errors:
        checks.append(
            "capability-capsule safety guidance is consistent across "
            f"{len(_CAPSULE_DOCUMENT_REQUIREMENTS)} policy and documentation files"
        )


def _validate_retention_prune_contract(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Keep the enabled retained-evidence expiration boundary documented."""

    starting_errors = len(errors)
    for relative, requirements in _RETENTION_PRUNE_DOCUMENT_REQUIREMENTS.items():
        path = root / relative
        try:
            text = path.read_text(encoding="utf-8")
        except (OSError, UnicodeError) as error:
            errors.append(
                f"retention-prune contract document is unreadable: {relative}: {error}"
            )
            continue
        for requirement in requirements:
            if requirement not in text:
                errors.append(
                    "retention-prune contract document is inconsistent: "
                    f"{relative} is missing {requirement!r}"
                )
        for forbidden in _RETENTION_PRUNE_FORBIDDEN_CURRENT_CLAIMS.get(relative, ()):
            if forbidden in text:
                errors.append(
                    "retention-prune contract contains a stale preview-only claim: "
                    f"{relative} contains {forbidden!r}"
                )

    if len(errors) == starting_errors:
        checks.append(
            "retained-evidence expiration guidance is consistent across "
            f"{len(_RETENTION_PRUNE_DOCUMENT_REQUIREMENTS)} documentation files"
        )


def _parse_agent_frontmatter(
    text: str,
) -> tuple[dict[str, object], str, tuple[str, ...]]:
    """Parse the deliberately small YAML subset used by the agent profile."""

    lines = text.splitlines()
    if not lines or lines[0] != "---":
        return {}, "", ("file must start with a YAML delimiter",)
    try:
        closing = lines.index("---", 1)
    except ValueError:
        return {}, "", ("closing YAML delimiter is missing",)

    metadata: dict[str, object] = {}
    parse_errors: list[str] = []
    active_list: str | None = None
    for number, line in enumerate(lines[1:closing], start=2):
        if line.startswith("  - "):
            if active_list is None:
                parse_errors.append(f"line {number} has a list item without a key")
                continue
            value = line[4:].strip()
            if not value:
                parse_errors.append(f"line {number} has an empty list item")
                continue
            existing = metadata.get(active_list)
            if not isinstance(existing, list):
                parse_errors.append(f"line {number} cannot extend {active_list}")
                continue
            existing.append(value)
            continue
        active_list = None
        match = _AGENT_FRONTMATTER_KEY.fullmatch(line)
        if match is None:
            parse_errors.append(f"line {number} is not a supported key or list item")
            continue
        key, raw_value = match.groups()
        if key in metadata:
            parse_errors.append(f"line {number} repeats key {key}")
            continue
        if not raw_value:
            metadata[key] = []
            active_list = key
        elif raw_value == "true":
            metadata[key] = True
        elif raw_value == "false":
            metadata[key] = False
        else:
            metadata[key] = raw_value.strip("'\"")

    for key, value in tuple(metadata.items()):
        if isinstance(value, list):
            metadata[key] = tuple(value)
    body = "\n".join(lines[closing + 1 :]).strip()
    if not body:
        parse_errors.append("instruction body is empty")
    return metadata, body, tuple(parse_errors)


def _validate_markdown_links(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    count = 0
    for path in _iter_files(root, suffixes={".md"}):
        text = path.read_text(encoding="utf-8")
        for target in _MARKDOWN_LINK.findall(text):
            clean = target.strip().split(maxsplit=1)[0].strip("<>")
            if not clean or clean.startswith(("#", "http://", "https://", "mailto:")):
                continue
            file_part = clean.split("#", 1)[0]
            resolved = (path.parent / file_part).resolve()
            try:
                resolved.relative_to(root)
            except ValueError:
                errors.append(f"Markdown link escapes project root: {path}: {target}")
                continue
            if not resolved.exists():
                errors.append(f"broken Markdown link: {path}: {target}")
            else:
                count += 1
    if not any(item.startswith("broken Markdown") for item in errors):
        checks.append(f"validated {count} local Markdown links")


def _validate_documentation(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Cross-check release documentation and checked-in plan examples."""

    readme = (root / "README.md").read_text(encoding="utf-8")
    changelog = (root / "CHANGELOG.md").read_text(encoding="utf-8")
    release_guide = (root / "docs/release-validation.md").read_text(encoding="utf-8")
    cli_reference = (root / "docs/cli-reference.md").read_text(encoding="utf-8")
    pyproject = tomllib.loads((root / "pyproject.toml").read_text(encoding="utf-8"))
    version = str(pyproject["project"]["version"])

    version_claims = {
        "README version banner": f"**Version {version} —",
        "release-validation title": f"# Release Validation — v{version}",
        "source-distribution filename": f"master_agent-{version}.tar.gz",
        "wheel filename": f"master_agent-{version}-py3-none-any.whl",
    }
    version_text = {
        "README version banner": readme,
        "release-validation title": release_guide,
        "source-distribution filename": readme,
        "wheel filename": readme,
    }
    for label, claim in version_claims.items():
        if claim not in version_text[label]:
            errors.append(f"documentation has stale {label}: expected {claim}")

    docs = {path.relative_to(root).as_posix() for path in (root / "docs").glob("*.md")}
    indexed = {
        target.split("#", 1)[0]
        for target in _MARKDOWN_LINK.findall(readme)
        if target.startswith("docs/")
    }
    missing_docs = sorted(docs - indexed)
    if missing_docs:
        errors.append(
            "README documentation index is missing: " + ", ".join(missing_docs)
        )

    catalog = tomllib.loads(
        (root / "config/capabilities.toml").read_text(encoding="utf-8")
    )["capabilities"]
    risk_counts: dict[str, int] = {}
    for definition in catalog.values():
        risk = str(definition["risk"])
        risk_counts[risk] = risk_counts.get(risk, 0) + 1
    readme_claims = (
        f"**{len(catalog)} typed capabilities**",
        f"{risk_counts.get('read_only', 0)} read-only capabilities",
        f"{risk_counts.get('local_generation', 0)} local-generation capabilities",
        f"{risk_counts.get('reversible_write', 0)} reversible-write definitions",
        f"{risk_counts.get('external_communication', 0)} external-communication capabilities",
        f"{risk_counts.get('high_impact', 0)} high-impact capability",
    )
    for claim in readme_claims:
        if claim not in readme:
            errors.append(f"README capability summary is stale: expected {claim}")

    current_changelog = changelog.split("\n## ", 2)[1]
    article = "an" if str(len(catalog))[0] in "8" else "a"
    changelog_claim = f"{article} {len(catalog)}-capability catalog"
    if not current_changelog.startswith(f"{version} "):
        errors.append(f"CHANGELOG first release section is not version {version}")
    elif changelog_claim not in current_changelog:
        errors.append(
            f"CHANGELOG capability summary is stale: expected {changelog_claim}"
        )

    models = (root / "src/master_agent/models.py").read_text(encoding="utf-8")
    schema_match = re.search(r'schema_version: str = "([^"]+)"', models)
    if schema_match is None:
        errors.append("could not determine the current ChangePlan schema version")
    else:
        current_schema = schema_match.group(1)
        for path in sorted((root / "examples").glob("*plan.json")):
            try:
                plan = json.loads(path.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError) as error:
                errors.append(f"example plan could not be read: {path}: {error}")
                continue
            if plan.get("schema_version") != current_schema:
                errors.append(
                    f"example plan uses stale schema {path}: expected {current_schema}"
                )

    cli = (root / "src/master_agent/cli.py").read_text(encoding="utf-8")
    commands = set(re.findall(r'subparsers\.add_parser\(\s*"([^"]+)"', cli))
    missing_commands = sorted(
        command for command in commands if f"`{command}`" not in cli_reference
    )
    if missing_commands:
        errors.append(
            "CLI reference is missing commands: " + ", ".join(missing_commands)
        )

    if not any(
        item.startswith(
            (
                "documentation has stale",
                "README documentation index",
                "README capability summary",
                "CHANGELOG",
                "could not determine the current ChangePlan",
                "example plan",
                "CLI reference",
            )
        )
        for item in errors
    ):
        checks.append(
            f"documentation matches version {version}, {len(docs)} guides, "
            f"{len(catalog)} capabilities, {len(commands)} CLI commands, and "
            "current plan schema"
        )


def _validate_demo(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    demo_root = root / "examples/v1-demo"
    manifest_path = demo_root / "manifest.json"
    if not manifest_path.exists():
        errors.append("v1 demonstration manifest is missing")
        return
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if manifest.get("published") is not False:
        errors.append("v1 demonstration must be marked unpublished")
    if manifest.get("uses_live_credentials") is not False:
        errors.append("v1 demonstration must be marked credential-free")
    entries = manifest.get("files")
    if not isinstance(entries, list) or not entries:
        errors.append("v1 demonstration manifest has no files")
        return
    verified = 0
    for entry in entries:
        if not isinstance(entry, dict):
            errors.append("v1 demonstration manifest entry must be an object")
            continue
        relative = Path(str(entry.get("path", "")))
        if relative.is_absolute() or ".." in relative.parts:
            errors.append(f"unsafe v1 demonstration path: {relative}")
            continue
        path = (demo_root / relative).resolve()
        try:
            path.relative_to(demo_root.resolve())
        except ValueError:
            errors.append(f"v1 demonstration path escapes root: {relative}")
            continue
        if not path.is_file():
            errors.append(f"v1 demonstration file is missing: {relative}")
            continue
        if path.stat().st_size != int(entry.get("size", -1)):
            errors.append(f"v1 demonstration size mismatch: {relative}")
            continue
        if sha256_file(path) != str(entry.get("sha256", "")):
            errors.append(f"v1 demonstration digest mismatch: {relative}")
            continue
        verified += 1
    if verified == len(entries):
        checks.append(f"validated {verified} v1 demonstration artifacts")

    _validate_demo_readiness(root, demo_root, checks, errors)

    draft_manifest_path = demo_root / "draft-package/manifest.json"
    draft_manifest = json.loads(draft_manifest_path.read_text(encoding="utf-8"))
    if draft_manifest.get("published") is not False:
        errors.append("draft package must be marked unpublished")
    for entry in draft_manifest.get("artifacts", []):
        relative = Path(str(entry["path"]))
        path = draft_manifest_path.parent / relative
        if not path.is_file() or sha256_file(path) != str(entry["sha256"]):
            errors.append(f"draft package artifact mismatch: {relative}")
    _validate_demo_powerpoint(demo_root, checks, errors)


def _validate_demo_powerpoint(
    demo_root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Validate the packaged demonstration deck with an installed reader."""

    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError:
        errors.append(
            "v1 demonstration PowerPoint validation requires the optional drafts extra"
        )
        return

    try:
        deck = Presentation(demo_root / "draft-package/change-package.pptx")
        if len(deck.slides) != 3:
            errors.append(
                f"v1 demonstration PowerPoint expected 3 slides, found {len(deck.slides)}"
            )
        else:
            checks.append("v1 demonstration PowerPoint opens with 3 slides")
    except (KeyError, OSError, PackageNotFoundError, TypeError, ValueError) as error:
        errors.append(f"v1 demonstration PowerPoint failed to open: {error}")


def _validate_demo_readiness(
    root: Path,
    demo_root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    """Cross-check demo readiness claims against the shipped capability catalog."""

    readiness_path = demo_root / "readiness.json"
    catalog_path = root / "config/capabilities.toml"
    try:
        readiness = json.loads(readiness_path.read_text(encoding="utf-8"))
        with catalog_path.open("rb") as handle:
            catalog = tomllib.load(handle)
    except (
        FileNotFoundError,
        json.JSONDecodeError,
        OSError,
        tomllib.TOMLDecodeError,
    ) as error:
        errors.append(f"v1 demonstration readiness could not be validated: {error}")
        return

    if not isinstance(readiness, dict):
        errors.append("v1 demonstration readiness must be an object")
        return
    raw_checks = readiness.get("checks")
    if not isinstance(raw_checks, list):
        errors.append("v1 demonstration readiness checks must be a list")
        return
    coverage = next(
        (
            item
            for item in raw_checks
            if isinstance(item, dict) and item.get("name") == "governance_coverage"
        ),
        None,
    )
    capabilities = catalog.get("capabilities")
    if coverage is None or not isinstance(capabilities, dict):
        errors.append("v1 demonstration readiness lacks governance coverage metadata")
        return
    observed = coverage.get("covered_capabilities")
    expected = len(capabilities)
    if not isinstance(observed, int) or isinstance(observed, bool):
        errors.append("v1 demonstration covered_capabilities must be an integer")
    elif observed != expected:
        errors.append(
            "v1 demonstration readiness capability count does not match the "
            f"catalog ({observed} != {expected})"
        )
    else:
        checks.append(f"v1 demonstration readiness covers all {expected} capabilities")


def _validate_file_hygiene(
    root: Path,
    checks: list[str],
    errors: list[str],
) -> None:
    file_count = 0
    runtime_directory = root / ".master-agent"
    if runtime_directory.exists():
        errors.append(
            f"forbidden runtime directory in source tree: {runtime_directory}"
        )
    for path in root.rglob("*"):
        if any(
            part in _IGNORED_DIRS or part.endswith(".egg-info") for part in path.parts
        ):
            continue
        if path.is_symlink():
            errors.append(f"source release must not contain symlink: {path}")
            continue
        if not path.is_file():
            continue
        file_count += 1
        forbidden_environment = (
            path.name.startswith(".env") and path.name != ".env.example"
        )
        if (
            path.name in _FORBIDDEN_NAMES
            or path.suffix.lower() in _FORBIDDEN_SUFFIXES
            or forbidden_environment
        ):
            errors.append(f"forbidden runtime/secret file in source tree: {path}")
    if not any(
        "forbidden runtime" in item
        or "forbidden runtime/secret" in item
        or "symlink" in item
        for item in errors
    ):
        checks.append(f"source-tree hygiene passed for {file_count} files")


def _validate_archive_member(name: str, errors: list[str]) -> None:
    """Reject unsafe paths and runtime/secret files in a built archive."""

    relative = PurePosixPath(name)
    if relative.is_absolute() or ".." in relative.parts:
        errors.append(f"unsafe path in release archive: {name}")
        return
    if ".master-agent" in relative.parts:
        errors.append(f"forbidden runtime directory in release archive: {name}")
    filename = relative.name
    suffix = Path(filename).suffix.lower()
    forbidden_environment = filename.startswith(".env") and filename != ".env.example"
    if (
        filename in _FORBIDDEN_NAMES
        or filename == ".DS_Store"
        or suffix in _FORBIDDEN_SUFFIXES
        or forbidden_environment
    ):
        errors.append(f"forbidden runtime/secret file in release archive: {name}")


def _validate_archive_runtime_mode(
    name: str,
    mode: int,
    *,
    is_regular: bool,
    errors: list[str],
) -> None:
    """Reject a capsule worker another operating-system account can alter."""

    worker_suffixes = (
        "master_agent/capsule_worker.py",
        "master_agent/platform_runtime/posix/capsule_worker.py",
    )
    if not name.endswith(worker_suffixes):
        return
    if not is_regular:
        errors.append(f"release archive capsule worker is not regular: {name}")
    permissions = stat.S_IMODE(mode)
    if permissions & (stat.S_IWGRP | stat.S_IWOTH):
        errors.append(
            "release archive capsule worker is writable by group or others: "
            f"{name} ({permissions:#05o})"
        )


def _consume_stream(handle: BinaryIO) -> None:
    """Read an archive member fully so integrity errors are observed."""

    while handle.read(1024 * 1024):
        pass


def _iter_files(root: Path, *, suffixes: set[str]) -> Iterable[Path]:
    for path in root.rglob("*"):
        if any(
            part in _IGNORED_DIRS or part.endswith(".egg-info") for part in path.parts
        ):
            continue
        if path.is_file() and path.suffix.lower() in suffixes:
            yield path


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--root", type=Path, default=Path.cwd())
    parser.add_argument(
        "--archive",
        action="append",
        type=Path,
        default=[],
        help="also validate a built wheel or source archive",
    )
    parser.add_argument("--output", type=Path)
    return parser


def main() -> int:
    """Run release validation and return a shell-compatible status."""

    args = _build_parser().parse_args()
    report = validate_project(args.root)
    checks = list(report.checks)
    errors = list(report.errors)
    for archive in args.archive:
        archive_report = validate_archive(archive)
        checks.extend(archive_report.checks)
        errors.extend(archive_report.errors)
    for check in checks:
        print(f"PASS {check}")
    for error in errors:
        print(f"FAIL {error}", file=sys.stderr)
    if args.output is not None:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(
            json.dumps(
                ValidationReport(tuple(checks), tuple(errors)).to_dict(),
                indent=2,
            )
            + "\n",
            encoding="utf-8",
        )
    return 0 if not errors else 1


if __name__ == "__main__":
    raise SystemExit(main())
